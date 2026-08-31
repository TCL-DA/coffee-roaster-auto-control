# -*- coding: utf-8 -*-
"""
Bài giảng TOÁN LỚP 3 — Bài 25: Phép chia hết, Phép chia có dư
(Kết nối tri thức, tập 1, trang 72–74 · Chủ đề 4: Phép nhân, phép chia trong phạm vi 100)

Trình tự theo Công văn 2345: Khởi động → Khám phá → Luyện tập → Vận dụng → Củng cố.

Nền SÁNG, chữ to, có hình minh hoạ — người xem là học sinh 8 tuổi ngồi dưới lớp,
khác hẳn hai bộ slide kỹ thuật nền tối.

Chạy:  python build_toan3_deck.py
Ra:    Toan3-Bai25-Phep-chia-het-chia-co-du.pptx
"""
import os

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from otl_deck import Deck

# ─────────────────────────────────────────────────────────────────────────────
# BẢNG MÀU — sáng, ấm, thân thiện với trẻ nhỏ
# ─────────────────────────────────────────────────────────────────────────────
PALETTE = {
    "bg":      "FFF9F0",   # kem ấm
    "bg2":     "FFFFFF",
    "card":    "FFFFFF",
    "cardwarm": "FFF3E0",
    "stroke":  "EFE0C9",
    "accent":  "FF6B35",   # cam — phép chia có dư
    "accent2": "1FA971",   # xanh lá — phép chia hết
    "blue":    "2E7FE8",
    "gold":    "FFB800",
    "grape":   "7B5CFF",
    "ink":     "241C14",   # nâu gần đen
    "mute":    "6B5E50",
    "dim":     "A0917F",
    "grid":    "EFE0C9",
    "leaf":    "3AAE5A",
    "stem":    "8A5A32",
    "plate":   "E8EDF2",
}

D = Deck(PALETTE, total=12)
M, CW, W, H = D.M, D.CW, D.W, D.H


def IN(v):
    return int(Inches(v))


# ─────────────────────────────────────────────────────────────────────────────
# HÌNH MINH HOẠ
# ─────────────────────────────────────────────────────────────────────────────
def apple(s, cx, cy, r, color="accent"):
    """Một quả táo: thân tròn + cuống + lá."""
    body = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - r), int(cy - r * 0.92),
                              int(2 * r), int(1.9 * r))
    D.solid(body, color)
    D._no_line(body)
    stem = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              int(cx - r * 0.07), int(cy - r * 1.35),
                              int(r * 0.15), int(r * 0.5))
    D.solid(stem, "stem")
    D._no_line(stem)
    leaf = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx + r * 0.05), int(cy - r * 1.3),
                              int(r * 0.55), int(r * 0.3))
    D.solid(leaf, "leaf")
    D._no_line(leaf)
    return body


def plate(s, cx, cy, w):
    """Cái đĩa nhìn nghiêng."""
    h = int(w * 0.34)
    p = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - w / 2), int(cy - h / 2), int(w), h)
    D.solid(p, "plate")
    p.line.color.rgb = __import__("pptx").dml.color.RGBColor.from_string(PALETTE["stroke"])
    p.line.width = Pt(1.5)
    return p


def long_division(s, left, top, dividend, divisor, quotient, product, remainder,
                  scale=1.0):
    """Vẽ khối đặt tính cột dọc:   9 | 2
                                  -8 | 4
                                   1
    """
    fs = 32 * scale
    colw = IN(0.62 * scale)
    rowh = IN(0.5 * scale)
    barx = left + colw * 2
    # số bị chia
    D.text(s, left, top, colw * 2, rowh, [(str(dividend), fs, "num", "ink")],
           align=PP_ALIGN.RIGHT)
    # vạch dọc
    v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, barx + IN(0.1), top,
                           IN(0.028), int(rowh * 2))
    D.solid(v, "dim")
    D._no_line(v)
    # số chia
    D.text(s, barx + IN(0.26), top, colw * 2, rowh, [(str(divisor), fs, "num", "ink")])
    # tích trừ đi
    D.text(s, left, top + rowh, colw * 2, rowh,
           [("−" + str(product), fs, "num", "mute")], align=PP_ALIGN.RIGHT)
    # thương
    D.text(s, barx + IN(0.26), top + rowh, colw * 2, rowh,
           [(str(quotient), fs, "num", "accent2")])
    # vạch ngang
    hbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + IN(0.1), top + rowh * 2,
                              barx - left, IN(0.028))
    D.solid(hbar, "dim")
    D._no_line(hbar)
    # số dư
    col = "accent" if remainder else "dim"
    D.text(s, left, top + int(rowh * 2.1), colw * 2, rowh,
           [(str(remainder), fs, "num", col)], align=PP_ALIGN.RIGHT)
    return top + int(rowh * 3.1)


def foot(s, n):
    """Chân trang của bài giảng — không dùng chân trang công ty."""
    D.text(s, M, H - IN(0.72), CW, IN(0.3),
           [("Toán 3  ·  Bài 25: Phép chia hết, phép chia có dư", 10.5, "body", "dim")])
    D.text(s, M, H - IN(0.72), CW, IN(0.3),
           [("%d / 12" % n, 10.5, "num", "dim", 1.2)], align=PP_ALIGN.RIGHT)


def stage(s, txt, color="blue"):
    """Nhãn giai đoạn tiết học ở góc phải trên."""
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             int(W - M - IN(2.6)), IN(0.66), IN(2.6), IN(0.44))
    tag.adjustments[0] = 0.45
    D.solid(tag, color)
    D._no_line(tag)
    D.text(s, int(W - M - IN(2.6)), IN(0.78), IN(2.6), IN(0.3),
           [(txt.upper(), 12, "bold", "bg2", 1.8)], align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 01 · BÌA
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Giới thiệu bài: Bài 25 Toán 3 Kết nối tri thức, trang 72–74.")
D.orb(s, IN(8.6), IN(-1.4), IN(7.6), "gold", 26)
D.orb(s, IN(-2.0), IN(4.4), IN(5.6), "accent2", 16, "ORB2")
D.text(s, M, IN(1.35), CW, IN(0.4),
       [("TOÁN LỚP 3  ·  KẾT NỐI TRI THỨC", 13, "bold", "blue", 2.4)], name="KICKER")
D.rule(s, M, IN(1.92), IN(1.5), "accent", Pt(5), "RULE")
D.text(s, M, IN(2.35), IN(11.6), IN(2.4),
       [("Phép chia hết", 62, "display", "accent2", -1.5),
        ("Phép chia có dư", 62, "display", "accent", -1.5)], line=1.12, name="HEADING")
D.text(s, M, IN(4.95), IN(9.6), IN(0.5),
       [("Bài 25  ·  Sách giáo khoa trang 72 – 74", 22, "light", "mute")])
D.rule(s, M, IN(5.7), IN(11.6), "stroke", Pt(1))
D.text(s, M, IN(5.98), IN(11.6), IN(0.4),
       [("Chủ đề 4: Phép nhân, phép chia trong phạm vi 100", 14, "body", "dim")])
for i in range(4):
    apple(s, IN(9.4 + i * 0.85), IN(5.2), IN(0.3),
          "accent" if i % 2 == 0 else "accent2")

# ─────────────────────────────────────────────────────────────────────────────
# 02 · YÊU CẦU CẦN ĐẠT
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Yêu cầu cần đạt của tiết học.")
D.orb(s, IN(10.2), IN(4.4), IN(6.0), "gold", 18)
D.kicker(s, "Sau bài học, em sẽ", color="blue")
D.heading(s, "Em học được gì hôm nay?")
goals = [
    ("Nhận ra phép chia hết", "Chia xong không thừa cái nào", "accent2"),
    ("Nhận ra phép chia có dư", "Chia xong vẫn còn thừa lại", "accent"),
    ("Nhớ quy tắc số dư", "Số dư luôn bé hơn số chia", "grape"),
    ("Giải bài toán thực tế", "Chia cá vào rổ, chia người lên thuyền", "blue"),
]
cw4 = (CW - IN(0.6)) / 4
for i, (t, d, col) in enumerate(goals):
    x = int(M + i * (cw4 + IN(0.2)))
    D.card(s, x, IN(2.8), int(cw4), IN(2.7))
    n = s.shapes.add_shape(MSO_SHAPE.OVAL, x + IN(0.4), IN(3.2), IN(0.62), IN(0.62))
    D.solid(n, col)
    D._no_line(n)
    D.text(s, x + IN(0.4), IN(3.34), IN(0.62), IN(0.4),
           [(str(i + 1), 24, "display", "bg2")], align=PP_ALIGN.CENTER)
    D.text(s, x + IN(0.4), IN(4.1), int(cw4 - IN(0.8)), IN(0.7),
           [(t, 19, "display", "ink", -0.3)], line=1.1)
    D.text(s, x + IN(0.4), IN(4.85), int(cw4 - IN(0.8)), IN(0.5),
           [(d, 13.5, "body", "mute")], line=1.3)
foot(s, 2)

# ─────────────────────────────────────────────────────────────────────────────
# 03 · KHỞI ĐỘNG
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Khởi động: 8 quả táo chia đều vào 2 đĩa thì mỗi đĩa mấy quả?")
D.orb(s, IN(-2.4), IN(-1.6), IN(6.4), "gold", 22)
D.kicker(s, "Bài toán mở đầu", color="gold")
stage(s, "Khởi động", "gold")
D.heading(s, "Cô có 8 quả táo và 2 cái đĩa")
D.text(s, M, IN(2.35), IN(11.0), IN(0.6),
       [("Chia đều số táo vào hai đĩa thì mỗi đĩa được mấy quả?", 26, "light", "mute")])
for i in range(8):
    apple(s, IN(1.5 + i * 0.95), IN(3.9), IN(0.36))
plate(s, IN(4.0), IN(5.5), IN(3.0))
plate(s, IN(9.2), IN(5.5), IN(3.0))
D.text(s, IN(2.5), IN(6.22), IN(3.0), IN(0.4),
       [("Đĩa 1", 16, "bold", "mute")], align=PP_ALIGN.CENTER)
D.text(s, IN(7.7), IN(6.22), IN(3.0), IN(0.4),
       [("Đĩa 2", 16, "bold", "mute")], align=PP_ALIGN.CENTER)
foot(s, 3)

# ─────────────────────────────────────────────────────────────────────────────
# 04 · KHÁM PHÁ 1 — PHÉP CHIA HẾT
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Khám phá 1: 8 : 2 = 4, không thừa quả nào — phép chia hết.")
D.orb(s, IN(9.6), IN(-2.0), IN(6.8), "accent2", 18)
D.kicker(s, "Khám phá 1", color="accent2")
stage(s, "Khám phá", "accent2")
D.heading(s, "Mỗi đĩa 4 quả — không thừa quả nào")
plate(s, IN(3.0), IN(3.9), IN(3.4))
plate(s, IN(3.0), IN(5.3), IN(3.4))
for r in range(2):
    for i in range(4):
        apple(s, IN(1.85 + i * 0.78), IN(3.72 + r * 1.4), IN(0.3), "accent2")
D.text(s, IN(1.3), IN(6.05), IN(3.4), IN(0.4),
       [("8 quả  ·  2 đĩa  ·  mỗi đĩa 4 quả", 14, "body", "mute")], align=PP_ALIGN.CENTER)

D.card(s, IN(5.35), IN(2.85), IN(3.5), IN(3.3), "cardwarm", "stroke")
D.text(s, IN(5.35), IN(3.15), IN(3.5), IN(0.4),
       [("ĐẶT TÍNH", 12, "bold", "dim", 1.8)], align=PP_ALIGN.CENTER)
long_division(s, IN(6.15), IN(3.7), 8, 2, 4, 8, 0)
D.text(s, IN(5.35), IN(5.5), IN(3.5), IN(0.4),
       [("số dư bằng 0", 15, "bold", "accent2")], align=PP_ALIGN.CENTER)

D.card(s, IN(9.15), IN(2.85), IN(3.33), IN(3.3), "card", "accent2")
D.text(s, IN(9.55), IN(3.25), IN(2.6), IN(0.5),
       [("8 : 2 = 4", 34, "display", "ink", -1)])
D.rule(s, IN(9.55), IN(4.0), IN(0.7), "accent2", Pt(4))
D.text(s, IN(9.55), IN(4.25), IN(2.6), IN(1.6),
       [("Chia xong không còn thừa quả nào.", 15, "body", "mute")], line=1.35)
D.text(s, IN(9.55), IN(5.35), IN(2.6), IN(0.5),
       [("PHÉP CHIA HẾT", 17, "display", "accent2", 0.6)])
foot(s, 4)

# ─────────────────────────────────────────────────────────────────────────────
# 05 · KHÁM PHÁ 2 — PHÉP CHIA CÓ DƯ
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Khám phá 2: 9 : 2 = 4 dư 1 — phép chia có dư.")
D.orb(s, IN(9.6), IN(-2.0), IN(6.8), "accent", 20)
D.kicker(s, "Khám phá 2", color="accent")
stage(s, "Khám phá", "accent")
D.heading(s, "Thêm 1 quả nữa — thừa ra 1 quả")
plate(s, IN(3.0), IN(3.9), IN(3.4))
plate(s, IN(3.0), IN(5.3), IN(3.4))
for r in range(2):
    for i in range(4):
        apple(s, IN(1.85 + i * 0.78), IN(3.72 + r * 1.4), IN(0.3), "accent2")
# quả thừa
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(1.35), IN(5.95), IN(3.3), IN(0.7))
box.adjustments[0] = 0.25
D.solid(box, "cardwarm")
box.line.color.rgb = __import__("pptx").dml.color.RGBColor.from_string(PALETTE["accent"])
box.line.width = Pt(2)
apple(s, IN(1.95), IN(6.32), IN(0.26), "accent")
D.text(s, IN(2.35), IN(6.16), IN(2.2), IN(0.4),
       [("thừa 1 quả", 16, "bold", "accent")])

D.card(s, IN(5.35), IN(2.85), IN(3.5), IN(3.3), "cardwarm", "stroke")
D.text(s, IN(5.35), IN(3.15), IN(3.5), IN(0.4),
       [("ĐẶT TÍNH", 12, "bold", "dim", 1.8)], align=PP_ALIGN.CENTER)
long_division(s, IN(6.15), IN(3.7), 9, 2, 4, 8, 1)
D.text(s, IN(5.35), IN(5.5), IN(3.5), IN(0.4),
       [("số dư bằng 1", 15, "bold", "accent")], align=PP_ALIGN.CENTER)

D.card(s, IN(9.15), IN(2.85), IN(3.33), IN(3.3), "card", "accent")
D.text(s, IN(9.55), IN(3.25), IN(2.8), IN(0.5),
       [("9 : 2 = 4 (dư 1)", 26, "display", "ink", -1)])
D.rule(s, IN(9.55), IN(4.0), IN(0.7), "accent", Pt(4))
D.text(s, IN(9.55), IN(4.25), IN(2.6), IN(1.6),
       [("Mỗi đĩa vẫn 4 quả, nhưng còn thừa lại 1 quả.", 15, "body", "mute")], line=1.35)
D.text(s, IN(9.55), IN(5.35), IN(2.8), IN(0.5),
       [("PHÉP CHIA CÓ DƯ", 17, "display", "accent", 0.6)])
foot(s, 5)

# ─────────────────────────────────────────────────────────────────────────────
# 06 · SO SÁNH
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("So sánh hai phép chia cạnh nhau.")
D.orb(s, IN(5.4), IN(-2.2), IN(6.4), "gold", 16)
D.kicker(s, "So sánh", color="blue")
D.heading(s, "Hai phép chia khác nhau chỗ nào?")
half = (CW - IN(0.4)) / 2
cols = [(M, "8 : 2 = 4", "PHÉP CHIA HẾT", "accent2",
         ["Chia hết, không thừa gì", "Số dư bằng 0", "Viết gọn: 8 : 2 = 4"]),
        (int(M + half + IN(0.4)), "9 : 2 = 4 (dư 1)", "PHÉP CHIA CÓ DƯ", "accent",
         ["Chia xong vẫn thừa lại", "Số dư khác 0", "Viết: 9 : 2 = 4 (dư 1)"])]
for x, expr, title, col, items in cols:
    D.card(s, x, IN(2.75), int(half), IN(3.5), "card", col)
    D.text(s, x + IN(0.5), IN(3.1), int(half - IN(1.0)), IN(0.7),
           [(expr, 40, "display", "ink", -1.2)])
    D.text(s, x + IN(0.5), IN(3.95), int(half - IN(1.0)), IN(0.4),
           [(title, 15, "display", col, 1.2)])
    D.rule(s, x + IN(0.5), IN(4.45), IN(0.7), col, Pt(4))
    for j, it in enumerate(items):
        y = IN(4.7) + int(IN(0.42) * j)
        dotm = s.shapes.add_shape(MSO_SHAPE.OVAL, x + IN(0.5), y + IN(0.08),
                                  IN(0.14), IN(0.14))
        D.solid(dotm, col)
        D._no_line(dotm)
        D.text(s, x + IN(0.85), y, int(half - IN(1.4)), IN(0.35),
               [(it, 16, "body", "mute")])
foot(s, 6)

# ─────────────────────────────────────────────────────────────────────────────
# 07 · QUY TẮC VÀNG
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Quy tắc phải nhớ: số dư luôn bé hơn số chia.")
D.orb(s, IN(3.4), IN(0.6), IN(7.4), "gold", 26)
D.kicker(s, "Ghi nhớ", color="grape")
D.text(s, M, IN(2.1), IN(11.6), IN(1.6),
       [("Số dư luôn bé hơn số chia", 54, "display", "ink", -1.5)],
       line=1.05, name="HEADING")
D.rule(s, M, IN(3.35), IN(2.2), "grape", Pt(6))
ex = [("9 : 2 = 4 (dư 1)", "1 bé hơn 2", True),
      ("19 : 5 = 3 (dư 4)", "4 bé hơn 5", True),
      ("19 : 5 = 2 (dư 9)", "9 lớn hơn 5 — sai rồi!", False)]
ew = (CW - IN(0.5)) / 3
for i, (expr, note, ok) in enumerate(ex):
    x = int(M + i * (ew + IN(0.25)))
    col = "accent2" if ok else "accent"
    D.card(s, x, IN(3.9), int(ew), IN(1.9), "card", col)
    D.text(s, x + IN(0.4), IN(4.25), int(ew - IN(0.8)), IN(0.6),
           [(expr, 27, "display", "ink", -0.8)])
    mark = s.shapes.add_shape(MSO_SHAPE.OVAL, x + IN(0.4), IN(5.0), IN(0.36), IN(0.36))
    D.solid(mark, col)
    D._no_line(mark)
    D.text(s, x + IN(0.4), IN(5.05), IN(0.36), IN(0.3),
           [("✓" if ok else "✗", 16, "bold", "bg2")], align=PP_ALIGN.CENTER)
    D.text(s, x + IN(0.92), IN(5.06), int(ew - IN(1.3)), IN(0.35),
           [(note, 15, "bold", col)])
D.text(s, M, IN(6.15), CW, IN(0.4),
       [("Số bị chia  =  thương  ×  số chia  +  số dư", 18, "bold", "mute")])
foot(s, 7)

# ─────────────────────────────────────────────────────────────────────────────
# 08 · LUYỆN TẬP — NỐI
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Luyện tập bài 2 trang 74: nối phép chia với số dư của nó.")
D.orb(s, IN(10.4), IN(4.4), IN(5.8), "blue", 14)
D.kicker(s, "SGK trang 74  ·  Bài 2", color="blue")
stage(s, "Luyện tập", "blue")
D.heading(s, "Nối phép chia với số dư của nó")
quiz = ["17 : 2", "41 : 6", "19 : 7", "16 : 6", "34 : 6", "19 : 5"]
qw = (CW - IN(0.75)) / 6
for i, q in enumerate(quiz):
    x = int(M + i * (qw + IN(0.15)))
    D.card(s, x, IN(2.85), int(qw), IN(1.3), "card", "stroke")
    D.text(s, x, IN(3.25), int(qw), IN(0.5),
           [(q, 28, "display", "ink", -0.8)], align=PP_ALIGN.CENTER)
D.text(s, M, IN(4.45), CW, IN(0.4),
       [("SỐ DƯ", 13, "bold", "dim", 2.2)], align=PP_ALIGN.CENTER)
opts = ["1", "4", "5"]
ow = IN(1.5)
for i, o in enumerate(opts):
    x = int(M + (CW - (3 * ow + IN(0.6))) / 2 + i * (ow + IN(0.3)))
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, IN(4.95), ow, IN(1.1))
    c.adjustments[0] = 0.25
    D.solid(c, "cardwarm")
    c.line.color.rgb = __import__("pptx").dml.color.RGBColor.from_string(PALETTE["gold"])
    c.line.width = Pt(2)
    D.text(s, x, IN(5.25), ow, IN(0.5), [(o, 34, "display", "gold")],
           align=PP_ALIGN.CENTER)
D.text(s, M, IN(6.3), CW, IN(0.4),
       [("Gợi ý: chia thử rồi xem còn thừa mấy đơn vị.", 15, "body", "mute")],
       align=PP_ALIGN.CENTER)
foot(s, 8)

# ─────────────────────────────────────────────────────────────────────────────
# 09 · ĐÁP ÁN
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Đáp án bài nối.")
D.orb(s, IN(-2.2), IN(4.2), IN(5.6), "accent2", 16)
D.kicker(s, "SGK trang 74  ·  Bài 2", color="accent2")
stage(s, "Đáp án", "accent2")
D.heading(s, "Cùng kiểm tra nào!")
ans = [("17 : 2", "8", "1"), ("41 : 6", "6", "5"), ("19 : 7", "2", "5"),
       ("16 : 6", "2", "4"), ("34 : 6", "5", "4"), ("19 : 5", "3", "4")]
aw = (CW - IN(0.75)) / 6
for i, (q, th, du) in enumerate(ans):
    x = int(M + i * (aw + IN(0.15)))
    D.card(s, x, IN(2.85), int(aw), IN(2.5), "card", "accent2")
    D.text(s, x, IN(3.15), int(aw), IN(0.45),
           [(q, 24, "display", "ink", -0.8)], align=PP_ALIGN.CENTER)
    D.rule(s, x + int(aw * 0.3), IN(3.78), int(aw * 0.4), "stroke", Pt(2))
    D.text(s, x, IN(3.98), int(aw), IN(0.4),
           [("thương " + th, 14, "body", "mute")], align=PP_ALIGN.CENTER)
    D.text(s, x, IN(4.4), int(aw), IN(0.55),
           [("dư " + du, 26, "display", "accent")], align=PP_ALIGN.CENTER)
D.text(s, M, IN(5.75), CW, IN(0.5),
       [("Để ý: số dư ở cả sáu phép chia đều bé hơn số chia.", 18, "bold", "mute")],
       align=PP_ALIGN.CENTER)
foot(s, 9)

# ─────────────────────────────────────────────────────────────────────────────
# 10 · VẬN DỤNG 1 — CHIA CÁ VÀO RỔ
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Vận dụng bài 3 trang 74: 56 con cá, mỗi rổ 8 con.")
D.orb(s, IN(9.8), IN(-2.2), IN(6.4), "blue", 16)
D.kicker(s, "SGK trang 74  ·  Bài 3", color="blue")
stage(s, "Vận dụng", "blue")
D.heading(s, "Có 56 con cá, mỗi rổ đựng 8 con")
D.text(s, M, IN(2.4), IN(9.0), IN(0.5),
       [("Hỏi xếp được tất cả bao nhiêu rổ cá?", 24, "light", "mute")])
for i in range(7):
    x = int(M + i * IN(1.68))
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, IN(3.3), IN(1.45), IN(1.25))
    b.adjustments[0] = 0.2
    D.solid(b, "cardwarm")
    b.line.color.rgb = __import__("pptx").dml.color.RGBColor.from_string(PALETTE["blue"])
    b.line.width = Pt(2)
    D.text(s, x, IN(3.52), IN(1.45), IN(0.5),
           [("8", 30, "display", "blue")], align=PP_ALIGN.CENTER)
    D.text(s, x, IN(4.02), IN(1.45), IN(0.3),
           [("con cá", 12, "body", "mute")], align=PP_ALIGN.CENTER)
    D.text(s, x, IN(4.62), IN(1.45), IN(0.3),
           [("rổ %d" % (i + 1), 12, "bold", "dim")], align=PP_ALIGN.CENTER)
D.card(s, M, IN(5.35), IN(6.4), IN(1.25), "card", "accent2")
D.text(s, M + IN(0.5), IN(5.62), IN(5.4), IN(0.6),
       [("56 : 8 = 7 (rổ)", 34, "display", "ink", -1)])
D.card(s, int(M + IN(6.7)), IN(5.35), int(CW - IN(6.7)), IN(1.25), "cardwarm", "stroke")
D.text(s, int(M + IN(7.15)), IN(5.6), IN(4.2), IN(0.7),
       [("Chia hết — không thừa con cá nào.", 18, "body", "mute")], line=1.3)
foot(s, 10)

# ─────────────────────────────────────────────────────────────────────────────
# 11 · VẬN DỤNG 2 — THUYỀN
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Vận dụng nâng cao: 14 người, thuyền chở 4 người — cần mấy chuyến?")
D.orb(s, IN(-2.6), IN(-1.4), IN(6.2), "grape", 14)
D.kicker(s, "Bài toán thêm", color="grape")
stage(s, "Vận dụng", "grape")
D.heading(s, "14 người, mỗi thuyền chở được 4 người")
D.text(s, M, IN(2.4), IN(9.6), IN(0.5),
       [("Cần ít nhất bao nhiêu chuyến để chở hết mọi người?", 24, "light", "mute")])
for i in range(4):
    x = int(M + i * IN(3.0))
    full = i < 3
    col = "blue" if full else "accent"
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, IN(3.3), IN(2.6), IN(1.5))
    b.adjustments[0] = 0.18
    D.solid(b, "card")
    b.line.color.rgb = __import__("pptx").dml.color.RGBColor.from_string(PALETTE[col])
    b.line.width = Pt(2)
    for k in range(4 if full else 2):
        p = s.shapes.add_shape(MSO_SHAPE.OVAL, x + IN(0.35 + k * 0.52), IN(3.68),
                               IN(0.36), IN(0.36))
        D.solid(p, col)
        D._no_line(p)
    D.text(s, x, IN(4.25), IN(2.6), IN(0.35),
           [("chuyến %d  ·  %d người" % (i + 1, 4 if full else 2),
             13.5, "bold", "mute")], align=PP_ALIGN.CENTER)
D.card(s, M, IN(5.2), IN(5.6), IN(1.35), "card", "accent")
D.text(s, M + IN(0.5), IN(5.5), IN(4.6), IN(0.6),
       [("14 : 4 = 3 (dư 2)", 32, "display", "ink", -1)])
D.card(s, int(M + IN(5.9)), IN(5.2), int(CW - IN(5.9)), IN(1.35), "cardwarm", "accent")
D.text(s, int(M + IN(6.35)), IN(5.42), IN(5.0), IN(0.9),
       [("Còn thừa 2 người nên phải thêm một chuyến nữa.", 16, "body", "mute"),
        ("Vậy cần 4 chuyến.", 19, "display", "accent")], line=1.3)
foot(s, 11)

# ─────────────────────────────────────────────────────────────────────────────
# 12 · CỦNG CỐ
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Củng cố: ba điều cần nhớ.")
D.orb(s, IN(3.0), IN(0.5), IN(8.0), "gold", 24)
D.orb(s, IN(10.6), IN(4.6), IN(5.0), "accent2", 14, "ORB2")
D.text(s, M, IN(1.35), CW, IN(0.4),
       [("CỦNG CỐ BÀI HỌC", 13, "bold", "blue", 2.4)], name="KICKER")
D.rule(s, M, IN(1.92), IN(1.5), "accent", Pt(5), "RULE")
D.text(s, M, IN(2.3), IN(11.6), IN(1.4),
       [("Ba điều cần nhớ", 52, "display", "ink", -1.5)], name="HEADING")
rem = [("Chia hết", "Chia xong không thừa gì cả", "accent2"),
       ("Chia có dư", "Chia xong vẫn còn thừa lại", "accent"),
       ("Số dư", "Bao giờ cũng bé hơn số chia", "grape")]
rw = (CW - IN(0.5)) / 3
for i, (t, d, col) in enumerate(rem):
    x = int(M + i * (rw + IN(0.25)))
    D.card(s, x, IN(3.85), int(rw), IN(1.75), "card", col)
    D.rule(s, x + IN(0.45), IN(4.2), IN(0.6), col, Pt(4))
    D.text(s, x + IN(0.45), IN(4.45), int(rw - IN(0.9)), IN(0.5),
           [(t, 24, "display", "ink", -0.5)])
    D.text(s, x + IN(0.45), IN(5.0), int(rw - IN(0.9)), IN(0.4),
           [(d, 15, "body", "mute")])
D.text(s, M, IN(6.0), CW, IN(0.5),
       [("Về nhà làm bài 1, 2, 3 trang 74 nhé!", 22, "display", "blue")],
       align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Toan3-Bai25-Phep-chia-het-chia-co-du.pptx")
D.save(out)
print("Da luu:", out)
