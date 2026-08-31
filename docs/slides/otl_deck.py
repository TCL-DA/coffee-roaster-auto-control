# -*- coding: utf-8 -*-
"""
otl_deck.py — Bộ công cụ dựng slide O Tesla (nền tối, chuyển cảnh Morph).

Dùng chung cho mọi bộ slide OTL. Đổi bảng màu là đổi cả bộ.
Xem docs/slides/build_vacuum_deck.py để biết cách dùng.

QUY TẮC ĐÃ TRẢ GIÁ — đừng sửa lại thành cách cũ:
  · Pt(4) đã là EMU, KHÔNG nhân 12700 nữa.
  · Khối tô tự viết phải chèn NGAY SAU prstGeom, sai chỗ là shape ra trắng.
  · Textbox phải TẮT autofit, không thì PowerPoint kéo chữ về toạ độ 0.
  · Dựng xong PHẢI xuất PNG soi lại bằng mắt.
"""
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

# Ước lượng bề rộng chữ để canh không tràn dòng (inch / ký tự / pt).
CHAR_W_DISPLAY = 0.0074   # Segoe UI Black
CHAR_W_BODY    = 0.0062   # Segoe UI thường


def fits(txt, size_pt, width_in, display=True):
    """True nếu chuỗi vừa một dòng trong bề rộng cho trước."""
    w = CHAR_W_DISPLAY if display else CHAR_W_BODY
    return len(txt) * w * size_pt <= width_in


class Deck:
    """Một bộ slide 16:9 nền tối."""

    def __init__(self, palette, fonts=None, total=0, margin=0.85):
        self.C = palette
        self.F = fonts or {
            "display": "Segoe UI Black",
            "bold": "Segoe UI Semibold",
            "body": "Segoe UI",
            "light": "Segoe UI Light",
            "num": "Bahnschrift SemiBold Condensed",
        }
        self.total = total
        self.prs = Presentation()
        self.W, self.H = Inches(13.333), Inches(7.5)
        self.prs.slide_width, self.prs.slide_height = self.W, self.H
        self.blank = self.prs.slide_layouts[6]
        self.M = Inches(margin)
        self.CW = self.W - 2 * self.M

    # ── XML cấp thấp ────────────────────────────────────────────────────────
    @staticmethod
    def _no_line(shape):
        shape.line.fill.background()

    @staticmethod
    def _set_fill_xml(shape, xml):
        spPr = shape._element.spPr
        for tag in ("a:solidFill", "a:noFill", "a:gradFill",
                    "a:blipFill", "a:pattFill", "a:grpFill"):
            old = spPr.find(qn(tag))
            if old is not None:
                spPr.remove(old)
        node = etree.fromstring(xml)
        # lxml: phần tử rỗng là falsy — phải so `is None`, không dùng `or`.
        geom = spPr.find(qn("a:prstGeom"))
        if geom is None:
            geom = spPr.find(qn("a:custGeom"))
        if geom is not None:
            geom.addnext(node)
        else:
            spPr.insert(0, node)

    def solid(self, shape, color, alpha=None):
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(self.C.get(color, color))
        if alpha is not None:
            srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
            srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))}))
        return shape

    def glow(self, shape, color, center_alpha=45):
        """Quầng sáng toả tròn — thay soft glow của Canva."""
        c = self.C.get(color, color)
        self._set_fill_xml(shape, (
            '<a:gradFill %s rotWithShape="1"><a:gsLst>'
            '<a:gs pos="0"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
            '<a:gs pos="55000"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
            '<a:gs pos="100000"><a:srgbClr val="%s"><a:alpha val="0"/></a:srgbClr></a:gs>'
            '</a:gsLst>'
            '<a:path path="circle"><a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path>'
            "</a:gradFill>" % (nsdecls("a"), c, center_alpha * 1000,
                               c, int(center_alpha * 350), c)))
        return shape

    def linear(self, shape, c1, c2, angle=5400000):
        self._set_fill_xml(shape, (
            '<a:gradFill %s rotWithShape="1"><a:gsLst>'
            '<a:gs pos="0"><a:srgbClr val="%s"/></a:gs>'
            '<a:gs pos="100000"><a:srgbClr val="%s"/></a:gs>'
            '</a:gsLst><a:lin ang="%d" scaled="0"/></a:gradFill>'
            % (nsdecls("a"), self.C.get(c1, c1), self.C.get(c2, c2), angle)))
        return shape

    @staticmethod
    def _morph(slide, dur=900):
        slide._element.append(etree.fromstring(
            '<mc:AlternateContent '
            'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
            'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
            'xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" '
            'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<mc:Choice Requires="p159"><p:transition spd="slow" p14:dur="%d">'
            '<p159:morph option="byObject"/></p:transition></mc:Choice>'
            '<mc:Fallback><p:transition spd="slow"><p:fade/></p:transition></mc:Fallback>'
            "</mc:AlternateContent>" % dur))

    # ── Khối dựng ───────────────────────────────────────────────────────────
    def slide(self, notes=""):
        s = self.prs.slides.add_slide(self.blank)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.W, self.H)
        self.linear(bg, "bg2", "bg", 2700000)
        self._no_line(bg)
        bg.name = "CHROME_BG"   # nền: để Morph lo, không gán hiệu ứng vào
        if notes:
            s.notes_slide.notes_text_frame.text = notes
        self._morph(s)
        return s

    def orb(self, s, left, top, size, color="accent", alpha=42, name="CHROME_ORB"):
        """Vật thể xuyên suốt — Morph sẽ bay nó từ slide này sang slide kia."""
        e = s.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        self._no_line(e)
        self.glow(e, color, alpha)
        e.name = name
        return e

    def text(self, s, left, top, width, height, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line=None, name=None):
        """runs = [(chuỗi, cỡ pt, khoá font, khoá màu, giãn chữ pt)] — mỗi phần tử 1 đoạn."""
        tb = s.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE          # BẮT BUỘC — xem đầu file
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        for i, it in enumerate(runs):
            body, size, font, color = it[0], it[1], it[2], it[3]
            spc = it[4] if len(it) > 4 else 0
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if line:
                p.line_spacing = line
            r = p.add_run()
            r.text = body
            r.font.size = Pt(size)
            r.font.name = self.F.get(font, font)
            r.font.color.rgb = RGBColor.from_string(self.C.get(color, color))
            if spc:
                r.font._rPr.set("spc", str(int(spc * 100)))
        if name:
            tb.name = name
        return tb

    def rule(self, s, left, top, width, color="accent", thick=Pt(4), name=None):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(int(thick)))
        self.solid(r, color)
        self._no_line(r)
        if name:
            r.name = name
        return r

    def card(self, s, left, top, width, height, fill="card", stroke="stroke", name=None):
        c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        c.adjustments[0] = 0.06
        self.solid(c, fill)
        c.line.color.rgb = RGBColor.from_string(self.C.get(stroke, stroke))
        c.line.width = Pt(1)
        if name:
            c.name = name
        return c

    def polyline(self, s, pts, color="accent", width=Pt(2.5), dash=None, close=False):
        """pts = [(x_emu, y_emu), ...] — vẽ đường gấp khúc thật, không phải ảnh."""
        ff = s.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
        ff.add_line_segments(pts[1:], close=close)
        sh = ff.convert_to_shape()
        sh.fill.background()
        sh.line.color.rgb = RGBColor.from_string(self.C.get(color, color))
        sh.line.width = width
        if dash:
            ln = sh.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": dash}))
        return sh

    def dot(self, s, cx, cy, d, color="accent"):
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - d / 2), int(cy - d / 2), d, d)
        self.solid(o, color)
        self._no_line(o)
        return o

    # ── Thành phần lặp lại của bộ slide ─────────────────────────────────────
    def kicker(self, s, txt, top=Inches(0.72), color="accent"):
        self.text(s, self.M, top, self.CW, Inches(0.3),
                  [(txt.upper(), 12, "bold", color, 2.4)], name="KICKER")
        self.rule(s, self.M, top + Inches(0.32), Inches(1.5), color, Pt(5), "RULE")

    def heading(self, s, txt, size=40, top=Inches(1.12), color="ink", width=None):
        return self.text(s, self.M, top, width or self.CW, Inches(1.4),
                         [(txt, size, "display", color)], line=0.94, name="HEADING")

    def footer(self, s, n):
        self.text(s, self.M, self.H - Inches(0.72), self.CW, Inches(0.3),
                  [("O TESLA  ·  OTL-06ALS", 10, "body", "dim", 1.2)])
        if self.total:
            self.text(s, self.M, self.H - Inches(0.72), self.CW, Inches(0.3),
                      [("%02d / %02d" % (n, self.total), 10, "num", "dim", 1.2)],
                      align=PP_ALIGN.RIGHT)

    def save(self, path):
        self.prs.save(path)
        return path
