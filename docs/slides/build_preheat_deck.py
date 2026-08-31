# -*- coding: utf-8 -*-
"""
Dựng bộ slide giới thiệu tính năng PREHEAT (làm nóng máy rang tự động) — OTL-06ALS.
Tone đỏ/đen, chuyển cảnh Morph (kiểu Keynote Magic Move / iPhone motion).

Chạy:  python build_preheat_deck.py
Ra:    OTL-Preheat-Gioi-thieu.pptx  (cùng thư mục)
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

# ─────────────────────────────────────────────────────────────────────────────
# HỆ MÀU + CHỮ  (design token — đổi ở đây là đổi cả bộ)
# ─────────────────────────────────────────────────────────────────────────────
BG      = "0D0D10"   # nền đen sâu
CARD    = "17171D"   # nền thẻ
STROKE  = "2C2C35"   # viền mảnh
RED     = "E8342A"   # đỏ OTL
RED_BR  = "FF453A"   # đỏ sáng (iOS red)
EMBER   = "FF6A00"   # cam than lửa
GOLD    = "F5B841"
INK     = "FFFFFF"
MUTE    = "A0A0A9"   # chữ phụ
DIM     = "6C6C77"   # chữ mờ

F_DISPLAY = "Segoe UI Black"
F_BOLD    = "Segoe UI Semibold"
F_BODY    = "Segoe UI"
F_LIGHT   = "Segoe UI Light"
F_NUM     = "Bahnschrift SemiBold Condensed"

W, H   = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.85)
CW     = W - 2 * MARGIN          # bề rộng vùng nội dung

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ─────────────────────────────────────────────────────────────────────────────
# TIỆN ÍCH XML
# ─────────────────────────────────────────────────────────────────────────────
def _no_line(shape):
    shape.line.fill.background()


def solid(shape, hexcolor, alpha=None):
    """Tô đặc, có thể kèm độ trong suốt (alpha 0..100 = % đục)."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hexcolor)
    if alpha is not None:
        srgb = shape.fill.fore_color._xFill.find(qn("a:srgbClr"))
        a = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
        srgb.append(a)


def _set_fill_xml(shape, xml):
    """Thay khối tô của shape bằng XML tự viết.
    BẮT BUỘC chèn NGAY SAU prstGeom — sai thứ tự là PowerPoint bỏ qua, shape ra trắng."""
    from lxml import etree
    spPr = shape._element.spPr
    for tag in ("a:solidFill", "a:noFill", "a:gradFill",
                "a:blipFill", "a:pattFill", "a:grpFill"):
        old = spPr.find(qn(tag))
        if old is not None:
            spPr.remove(old)
    node = etree.fromstring(xml)
    geom = spPr.find(qn("a:prstGeom"))
    if geom is None:
        geom = spPr.find(qn("a:custGeom"))
    if geom is not None:
        geom.addnext(node)
    else:
        spPr.insert(0, node)


def radial_glow(shape, hexcolor, center_alpha=45):
    """Quầng sáng toả tròn — thay cho 'soft glow' của Canva."""
    xml = (
        '<a:gradFill %s rotWithShape="1">'
        '  <a:gsLst>'
        '    <a:gs pos="0"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '    <a:gs pos="55000"><a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr></a:gs>'
        '    <a:gs pos="100000"><a:srgbClr val="%s"><a:alpha val="0"/></a:srgbClr></a:gs>'
        '  </a:gsLst>'
        '  <a:path path="circle"><a:fillToRect l="50000" t="50000" r="50000" b="50000"/></a:path>'
        "</a:gradFill>" % (nsdecls("a"), hexcolor, center_alpha * 1000,
                           hexcolor, int(center_alpha * 1000 * 0.35), hexcolor)
    )
    _set_fill_xml(shape, xml)


def linear_fill(shape, c1, c2, angle=5400000):
    xml = (
        '<a:gradFill %s rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="%s"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="%s"/></a:gs>'
        '</a:gsLst><a:lin ang="%d" scaled="0"/></a:gradFill>'
        % (nsdecls("a"), c1, c2, angle)
    )
    _set_fill_xml(shape, xml)


def morph(slide, dur=900):
    """Chuyển cảnh Morph — vật thể cùng TÊN ở 2 slide sẽ tự bay/giãn sang nhau."""
    from lxml import etree
    xml = (
        '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
        '   xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
        '   xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" '
        '   xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        '  <mc:Choice Requires="p159">'
        '    <p:transition spd="slow" p14:dur="%d">'
        '      <p159:morph option="byObject"/>'
        "    </p:transition>"
        "  </mc:Choice>"
        "  <mc:Fallback>"
        '    <p:transition spd="slow"><p:fade/></p:transition>'
        "  </mc:Fallback>"
        "</mc:AlternateContent>" % dur
    )
    slide._element.append(etree.fromstring(xml))


# ─────────────────────────────────────────────────────────────────────────────
# TIỆN ÍCH DỰNG KHỐI
# ─────────────────────────────────────────────────────────────────────────────
def new_slide(notes=""):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    linear_fill(bg, "121218", BG, 2700000)
    _no_line(bg)
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    morph(s)
    return s


def ember(slide, left, top, size, color=RED, alpha=42, name="EMBER"):
    """Quầng lửa — vật thể xuyên suốt, Morph sẽ bay nó từ slide này sang slide kia."""
    e = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    _no_line(e)
    radial_glow(e, color, alpha)
    e.name = name
    return e


def text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line=None, space_after=0, name=None):
    """runs = [(nội dung, cỡ chữ, font, màu, giãn chữ pt)] — mỗi phần tử 1 đoạn."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    # TẮT autofit: để bật, PowerPoint tự co giãn hộp lúc mở file và chữ nhảy khỏi vị trí đã đặt.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, item in enumerate(runs):
        content, size, font, color = item[0], item[1], item[2], item[3]
        spc = item[4] if len(item) > 4 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line:
            p.line_spacing = line
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = content
        r.font.size = Pt(size)
        r.font.name = font
        r.font.color.rgb = RGBColor.from_string(color)
        if spc:
            r.font._rPr.set("spc", str(int(spc * 100)))
    if name:
        tb.name = name
    return tb


def rule(slide, left, top, width, color=RED, thick=Pt(4), name=None):
    # Pt() đã trả về Length tính bằng EMU — KHÔNG nhân 12700 lần nữa (nhân là phình khổng lồ).
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(int(thick)))
    solid(r, color)
    _no_line(r)
    if name:
        r.name = name
    return r


def card(slide, left, top, width, height, fill=CARD, stroke=STROKE, name=None):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.adjustments[0] = 0.06
    solid(c, fill)
    c.line.color.rgb = RGBColor.from_string(stroke)
    c.line.width = Pt(1)
    if name:
        c.name = name
    return c


def kicker(slide, txt, top=Inches(0.72), color=RED):
    """Nhãn nhỏ in hoa, giãn chữ — dấu hiệu nhận diện của bộ slide."""
    return text(slide, MARGIN, top, CW, Inches(0.3),
                [(txt.upper(), 12, F_BOLD, color, 2.4)], name="KICKER")


def heading(slide, txt, size=44, top=Inches(1.12), color=INK, width=None):
    return text(slide, MARGIN, top, width or CW, Inches(1.3),
                [(txt, size, F_DISPLAY, color)], line=0.94, name="HEADING")


def page_no(slide, n, total):
    text(slide, MARGIN, H - Inches(0.72), CW, Inches(0.3),
         [("O TESLA  ·  OTL-06ALS", 10, F_BODY, DIM, 1.2)])
    text(slide, MARGIN, H - Inches(0.72), CW, Inches(0.3),
         [("%02d / %02d" % (n, total), 10, F_NUM, DIM, 1.2)], align=PP_ALIGN.RIGHT)


TOTAL = 11

# ─────────────────────────────────────────────────────────────────────────────
# 01 · BÌA
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Mở đầu: tính năng Preheat — máy tự làm nóng tới đúng nhiệt trước khi vào mẻ.")
ember(s, Inches(7.4), Inches(-1.6), Inches(9.4), RED, 50)
ember(s, Inches(-2.2), Inches(4.2), Inches(6.4), EMBER, 26, "EMBER2")
rule(s, MARGIN, Inches(2.06), Inches(1.5), RED, Pt(5), "RULE")
text(s, MARGIN, Inches(1.5), CW, Inches(0.4),
     [("TÍNH NĂNG MỚI  ·  FIRMWARE OTL-06ALS", 12, F_BOLD, RED, 2.4)], name="KICKER")
text(s, MARGIN, Inches(2.5), Inches(9.6), Inches(2.4),
     [("PREHEAT", 138, F_DISPLAY, INK, -3)], line=0.86, name="HEADING")
text(s, MARGIN, Inches(4.55), Inches(9.6), Inches(1.4),
     [("Máy rang tự làm nóng — một nút, không cần canh", 26, F_LIGHT, MUTE)], line=1.15)
rule(s, MARGIN, Inches(5.55), Inches(11.6), STROKE, Pt(1))
text(s, MARGIN, Inches(5.85), Inches(11.6), Inches(0.4),
     [("O TESLA  ·  Điều khiển & Tự động hoá máy rang cà phê", 13, F_BODY, DIM, 1)])

# ─────────────────────────────────────────────────────────────────────────────
# 02 · VẤN ĐỀ
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Ba nỗi khổ khi làm nóng máy bằng tay.")
ember(s, Inches(9.9), Inches(-2.4), Inches(7.0), RED, 34)
ember(s, Inches(-3.0), Inches(5.4), Inches(5.0), EMBER, 18, "EMBER2")
kicker(s, "Vấn đề")
rule(s, MARGIN, Inches(1.04), Inches(1.5), RED, Pt(5), "RULE")
heading(s, "Làm nóng bằng tay — nửa tiếng đứng canh", 40)
pains = [
    ("20–30′", "đứng canh mỗi ca", "Thợ phải ngồi cạnh bảng điều khiển, chỉnh gas và gió liên tục cho tới khi trống đủ nóng."),
    ("7 kiểu", "bảy người bảy phách", "Mỗi thợ một tay nghề. Cùng một máy, cùng một mức nhiệt, ra bảy đường nhiệt khác nhau."),
    ("+15°C", "vọt nhiệt, hỏng mẻ đầu", "Ghi nhận thực tế: hạ từ 215°C xuống 150°C bị vọt tới +15°C — mẻ đầu tiên trong ngày luôn phải bỏ."),
]
cw3 = (CW - Inches(0.5)) / 3
for i, (big, small, desc) in enumerate(pains):
    x = MARGIN + i * (cw3 + Inches(0.25))
    c = card(s, x, Inches(2.75), cw3, Inches(3.15))
    rule(s, x + Inches(0.4), Inches(3.15), Inches(0.55), RED, Pt(4))
    text(s, x + Inches(0.4), Inches(3.5), cw3 - Inches(0.8), Inches(1.0),
         [(big, 46, F_DISPLAY, RED_BR, -1)], line=0.95)
    text(s, x + Inches(0.4), Inches(4.28), cw3 - Inches(0.8), Inches(0.4),
         [(small.upper(), 11, F_BOLD, GOLD, 1.6)])
    text(s, x + Inches(0.4), Inches(4.72), cw3 - Inches(0.8), Inches(1.2),
         [(desc, 13, F_BODY, MUTE)], line=1.3)
page_no(s, 2, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 03 · GIẢI PHÁP (câu tuyên ngôn)
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Tuyên ngôn: bấm một nút, máy tự lo phần còn lại.")
ember(s, Inches(3.4), Inches(0.9), Inches(7.4), RED, 46)
kicker(s, "Giải pháp", Inches(0.9))
rule(s, MARGIN, Inches(1.22), Inches(1.5), RED, Pt(5), "RULE")
text(s, MARGIN, Inches(2.3), Inches(11.6), Inches(2.2),
     [("Bấm một nút.", 60, F_DISPLAY, INK, -2),
      ("Máy tự lo phần còn lại.", 60, F_DISPLAY, RED_BR, -2)],
     line=1.1, name="HEADING")
text(s, MARGIN, Inches(4.85), Inches(10.4), Inches(1.2),
     [("Preheat là bộ điều khiển chạy trong firmware — đọc BT, ET và tốc độ tăng nhiệt "
       "mỗi giây, tự bậc gas và gió theo 8 trạng thái, rồi giữ đúng nhiệt chờ nạp hạt.",
       17, F_BODY, MUTE)], line=1.4)
page_no(s, 3, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 04 · CON SỐ
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Bốn con số tóm gọn tính năng.")
ember(s, Inches(10.2), Inches(4.6), Inches(6.2), EMBER, 24)
kicker(s, "Tóm tắt")
rule(s, MARGIN, Inches(1.04), Inches(1.5), RED, Pt(5), "RULE")
heading(s, "Bốn con số cần nhớ", 40)
stats = [
    ("≤ 3", "°C", "Vọt nhiệt tối đa", "Mục tiêu bám sát tay nghề thợ giỏi nhất"),
    ("8", "trạng thái", "Máy trạng thái", "Từ nguội, mồi lửa, lên nhiệt tới giữ nhiệt"),
    ("12", "vùng", "Bảng tự học", "Mỗi vùng 25°C, lưu lại trong thẻ nhớ"),
    ("5", "giây", "Nhịp ghi log", "60 dòng CSV mỗi lần preheat để soi lại"),
]
cw4 = (CW - Inches(0.6)) / 4
for i, (num, unit, title, sub) in enumerate(stats):
    x = MARGIN + i * (cw4 + Inches(0.2))
    card(s, x, Inches(2.7), cw4, Inches(3.2))
    text(s, x + Inches(0.35), Inches(3.05), cw4 - Inches(0.7), Inches(1.2),
         [(num, 62, F_DISPLAY, INK, -2)], line=0.92)
    text(s, x + Inches(0.35), Inches(4.1), cw4 - Inches(0.7), Inches(0.35),
         [(unit.upper(), 12, F_BOLD, RED_BR, 1.8)])
    rule(s, x + Inches(0.35), Inches(4.62), Inches(0.5), STROKE, Pt(2))
    text(s, x + Inches(0.35), Inches(4.85), cw4 - Inches(0.7), Inches(0.4),
         [(title, 15, F_BOLD, INK)])
    text(s, x + Inches(0.35), Inches(5.2), cw4 - Inches(0.7), Inches(0.7),
         [(sub, 12, F_BODY, DIM)], line=1.3)
page_no(s, 4, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 05 · MÁY TRẠNG THÁI
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Chuỗi 6 trạng thái chính; 2 trạng thái phụ xử lý máy đang nóng.")
ember(s, Inches(-2.4), Inches(-1.4), Inches(6.4), RED, 30)
kicker(s, "Cách máy chạy")
rule(s, MARGIN, Inches(1.04), Inches(1.5), RED, Pt(5), "RULE")
heading(s, "Sáu bước, máy tự đi hết", 40)
steps = [
    ("IDLE", "Phân loại", "Nguội / Nóng / Cần hạ"),
    ("PRE-IGNITE", "Thổi sạch", "Xả hết khí tồn trong buồng đốt"),
    ("IGNITE", "Mồi lửa", "Có xác nhận lửa, quá giờ là đóng gas"),
    ("HEATING", "Lên nhiệt", "Bậc gas theo tốc độ tăng nhiệt"),
    ("HOLDING", "Giữ nhiệt", "Ghìm quanh mức đặt, chờ nạp"),
    ("PRECISION", "Tinh chỉnh", "PI nhẹ, khử lệch còn ±2°C"),
]
bw = (CW - Inches(0.75)) / 6
for i, (code, name_, desc) in enumerate(steps):
    x = MARGIN + i * (bw + Inches(0.15))
    accent = RED_BR if i in (2, 5) else STROKE
    c = card(s, x, Inches(2.72), bw, Inches(2.5), CARD, accent)
    rule(s, x, Inches(2.72), bw, RED if i in (2, 5) else "3A3A45", Pt(3))
    text(s, x + Inches(0.24), Inches(3.02), bw - Inches(0.48), Inches(0.3),
         [("%02d" % (i + 1), 13, F_NUM, DIM, 1.2)])
    text(s, x + Inches(0.24), Inches(3.38), bw - Inches(0.48), Inches(0.5),
         [(code, 15, F_DISPLAY, INK, -0.3)], line=1.0)
    text(s, x + Inches(0.24), Inches(3.92), bw - Inches(0.48), Inches(0.35),
         [(name_, 14, F_BOLD, RED_BR)])
    text(s, x + Inches(0.24), Inches(4.32), bw - Inches(0.48), Inches(0.8),
         [(desc, 11.5, F_BODY, MUTE)], line=1.28)
    if i < 5:
        text(s, x + bw, Inches(3.62), Inches(0.15), Inches(0.3),
             [("›", 20, F_BODY, DIM)], align=PP_ALIGN.CENTER)
text(s, MARGIN, Inches(5.55), CW, Inches(0.9),
     [("Hai trạng thái phụ:  COOLING và COAST  —  khi máy đang nóng hơn mức cần, "
       "firmware tắt gas, mở gió và để nhiệt trôi xuống đúng chỗ thay vì mồi lửa thêm.",
       14, F_BODY, MUTE)], line=1.35)
page_no(s, 5, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 06 · HAI PHA ĐIỀU KHIỂN
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Nhanh trước, chuẩn sau — chia 60/40.")
ember(s, Inches(5.6), Inches(-2.2), Inches(6.6), RED, 32)
kicker(s, "Triết lý điều khiển")
rule(s, MARGIN, Inches(1.04), Inches(1.5), RED, Pt(5), "RULE")
heading(s, "Nhanh trước, chuẩn sau", 40)
half = (CW - Inches(0.3)) / 2
blocks = [
    (MARGIN, "60%", "RAMP  ·  3 phút đầu", EMBER,
     ["Ưu tiên lên nhiệt nhanh", "Gas bậc theo khoảng cách còn lại",
      "Gió ghìm thấp khi còn xa mức đặt", "Chấp nhận sai lệch ≤ 5°C"]),
    (MARGIN + half + Inches(0.3), "40%", "PRECISION  ·  2 phút cuối", RED_BR,
     ["Yêu cầu dao động trong ±2°C", "PI nhẹ khử lệch tĩnh",
      "Phanh sớm theo đà trôi nhiệt dự báo", "Kết thúc đúng mức đặt ±1°C"]),
]
for x, pct, title, col, items in blocks:
    card(s, x, Inches(2.7), half, Inches(3.4))
    text(s, x + Inches(0.45), Inches(3.0), half - Inches(0.9), Inches(0.9),
         [(pct, 50, F_DISPLAY, col, -1.5)], line=0.95)
    text(s, x + Inches(0.45), Inches(3.85), half - Inches(0.9), Inches(0.35),
         [(title.upper(), 12.5, F_BOLD, INK, 1.4)])
    rule(s, x + Inches(0.45), Inches(4.28), Inches(0.5), col, Pt(3))
    for j, it in enumerate(items):
        text(s, x + Inches(0.45), Inches(4.5) + Inches(0.34) * j, Inches(0.16), Inches(0.3),
             [("—", 12, F_BODY, col)])
        text(s, x + Inches(0.75), Inches(4.5) + Inches(0.34) * j, half - Inches(1.2), Inches(0.3),
             [(it, 13.5, F_BODY, MUTE)])
page_no(s, 6, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 07 · CHUỖI TRỄ NHIỆT
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Vì sao khó: giữa vặn gas và thấy trống nóng có tới ~30 giây trễ.")
ember(s, Inches(9.6), Inches(4.4), Inches(6.6), EMBER, 26)
kicker(s, "Vì sao khó")
rule(s, MARGIN, Inches(1.04), Inches(1.5), RED, Pt(5), "RULE")
heading(s, "Vặn gas bây giờ — 30 giây nữa trống mới nóng", 40)
chain = [
    ("GAS %", "lệnh vào", "tức thì", DIM),
    ("ĐẦU ĐỐT", "ngọn lửa ổn định", "2–3 giây", GOLD),
    ("BUỒNG ĐỐT", "khí nóng · cảm biến ET", "5–10 giây", EMBER),
    ("TRỐNG RANG", "kim loại · cảm biến BT", "15–20 giây", RED_BR),
]
bw = (CW - Inches(1.65)) / 4
for i, (nm, sub, lag, col) in enumerate(chain):
    x = MARGIN + i * (bw + Inches(0.55))
    card(s, x, Inches(2.85), bw, Inches(1.95))
    rule(s, x, Inches(2.85), bw, col, Pt(3))
    text(s, x + Inches(0.3), Inches(3.2), bw - Inches(0.6), Inches(0.5),
         [(nm, 19, F_DISPLAY, INK, -0.3)])
    text(s, x + Inches(0.3), Inches(3.68), bw - Inches(0.6), Inches(0.6),
         [(sub, 12.5, F_BODY, MUTE)], line=1.3)
    text(s, x + Inches(0.3), Inches(4.32), bw - Inches(0.6), Inches(0.35),
         [(lag.upper(), 11.5, F_BOLD, col, 1.2)])
    if i < 3:
        text(s, x + bw, Inches(3.55), Inches(0.55), Inches(0.4),
             [("→", 22, F_BODY, DIM)], align=PP_ALIGN.CENTER)
text(s, MARGIN, Inches(5.25), CW, Inches(1.1),
     [("Ba tầng trễ nối tiếp nhau là lý do người mới luôn vặn gas quá tay. "
       "Preheat xử lý bằng thời gian chết 15 giây sau mỗi nấc gas — bậc xong thì "
       "chờ nghe máy trả lời rồi mới bậc tiếp, đúng cách thợ giỏi vẫn làm.",
       15, F_BODY, MUTE)], line=1.4)
page_no(s, 7, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 08 · TỰ HỌC
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Càng chạy càng chuẩn: bảng feedforward 12 vùng + 3 tham số thích nghi.")
ember(s, Inches(-2.6), Inches(3.4), Inches(6.8), RED, 30)
kicker(s, "Tự học")
rule(s, MARGIN, Inches(1.04), Inches(1.5), RED, Pt(5), "RULE")
heading(s, "Chạy càng nhiều, máy càng chuẩn", 40)
text(s, MARGIN, Inches(2.32), Inches(7.4), Inches(0.6),
     [("Sau mỗi lần preheat, firmware ghi lại mức gas thực tế đã giữ được nhiệt "
       "ở từng vùng và dùng cho lần sau.", 15, F_BODY, MUTE)], line=1.35)
zones = [("0–25", 0), ("25–50", 0), ("50–75", 1), ("75–100", 1), ("100–125", 2),
         ("125–150", 2), ("150–175", 3), ("175–200", 3), ("200–225", 4),
         ("225–250", 4), ("250–275", 3), ("275–300", 2)]
zw = (Inches(7.4) - Inches(0.55)) / 6
for i, (lab, lv) in enumerate(zones):
    col_i, row_i = i % 6, i // 6
    x = MARGIN + col_i * (zw + Inches(0.11))
    y = Inches(3.35) + row_i * Inches(0.92)
    shade = ["1C1C23", "3A1F1D", "5E2620", "9B2E24", "E8342A"][lv]
    z = card(s, x, y, zw, Inches(0.78), shade, STROKE)
    text(s, x, y + Inches(0.14), zw, Inches(0.3),
         [(lab, 12.5, F_NUM, INK if lv >= 3 else MUTE)], align=PP_ALIGN.CENTER)
    text(s, x, y + Inches(0.42), zw, Inches(0.25),
         [("°C", 9.5, F_BODY, DIM, 1)], align=PP_ALIGN.CENTER)
text(s, MARGIN, Inches(5.28), Inches(7.4), Inches(0.4),
     [("12 VÙNG × 25°C  ·  LƯU TRONG /ph_ff.txt", 11, F_BOLD, DIM, 1.6)])

lx = MARGIN + Inches(7.85)
lw = CW - Inches(7.85)
card(s, lx, Inches(2.7), lw, Inches(3.4))
text(s, lx + Inches(0.45), Inches(3.0), lw - Inches(0.9), Inches(0.4),
     [("BA THAM SỐ MÁY TỰ ĐO", 11.5, F_BOLD, RED_BR, 1.6)])
adapt = [
    ("Đà trôi nhiệt", "Tắt gas rồi nhiệt còn trôi bao nhiêu"),
    ("Nhiệt tổn thất", "Máy nguội bao nhanh khi gas về nền"),
    ("Độ nhạy gas", "Thêm 1% gas thì nhiệt lên bao nhiêu"),
]
for j, (t, d) in enumerate(adapt):
    y = Inches(3.42) + Inches(0.85) * j
    rule(s, lx + Inches(0.4), y + Inches(0.06), Inches(0.28), RED, Pt(3))
    text(s, lx + Inches(0.4), y + Inches(0.2), lw - Inches(0.7), Inches(0.3),
         [(t, 15.5, F_BOLD, INK)])
    text(s, lx + Inches(0.4), y + Inches(0.5), lw - Inches(0.7), Inches(0.28),
         [(d, 11.5, F_BODY, MUTE)])
page_no(s, 8, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 09 · AN TOÀN
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Năm cờ lỗi tự giám sát + hai chốt cứng đóng gas.")
ember(s, Inches(9.4), Inches(-2.0), Inches(6.4), RED, 34)
kicker(s, "An toàn")
rule(s, MARGIN, Inches(1.04), Inches(1.5), RED, Pt(5), "RULE")
heading(s, "Máy tự canh chính nó", 40)
faults = [
    ("Mất tín hiệu cảm biến", "BT hoặc ET rơi đột ngột — nghi dây đo lỏng"),
    ("Tốc độ nhiệt bất thường", "RoR vọt ngoài dải hợp lý"),
    ("Mất lửa giữa chừng", "Đang đốt mà nhiệt không lên — nghi tắt lửa"),
    ("Mất liên lạc màn hình", "HMI không trả lời trong thời gian cho phép"),
    ("Lỗi thẻ nhớ", "Không ghi được log — vẫn chạy, chỉ mất dữ liệu"),
]
fw = (CW - Inches(0.4)) / 5
for i, (t, d) in enumerate(faults):
    x = MARGIN + i * (fw + Inches(0.1))
    card(s, x, Inches(2.72), fw, Inches(1.95))
    rule(s, x, Inches(2.72), fw, GOLD, Pt(3))
    text(s, x + Inches(0.26), Inches(3.02), fw - Inches(0.52), Inches(0.6),
         [(t, 14.5, F_BOLD, INK)], line=1.15)
    text(s, x + Inches(0.26), Inches(3.72), fw - Inches(0.52), Inches(0.8),
         [(d, 11.5, F_BODY, MUTE)], line=1.28)
hard = [
    ("Chênh ET − BT vượt 160°C", "Huỷ preheat, đóng gas ngay và báo lỗi lên màn hình."),
    ("Mồi lửa quá giờ cho phép", "Không xác nhận được lửa thì đóng gas, không thử vô hạn."),
]
for i, (t, d) in enumerate(hard):
    x = MARGIN + i * ((CW - Inches(0.3)) / 2 + Inches(0.3))
    ww = (CW - Inches(0.3)) / 2
    card(s, x, Inches(4.95), ww, Inches(1.2), "2A100E", RED)
    text(s, x + Inches(0.4), Inches(5.18), ww - Inches(0.8), Inches(0.35),
         [("CHỐT CỨNG  ·  " + t.upper(), 12, F_BOLD, RED_BR, 1.2)])
    text(s, x + Inches(0.4), Inches(5.55), ww - Inches(0.8), Inches(0.4),
         [(d, 13, F_BODY, MUTE)])
page_no(s, 9, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 10 · KẾT QUẢ
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Đối chiếu thợ làm tay và mục tiêu của Preheat.")
ember(s, Inches(10.4), Inches(4.2), Inches(6.0), EMBER, 22)
kicker(s, "Đối chiếu")
rule(s, MARGIN, Inches(1.04), Inches(1.5), RED, Pt(5), "RULE")
heading(s, "Thợ làm tay và máy tự làm", 40)
rows = [
    ("", "THỢ LÀM TAY", "PREHEAT TỰ ĐỘNG"),
    ("Người đứng canh", "Suốt 20–30 phút", "Không cần"),
    ("Vọt nhiệt điển hình", "+1 … +3°C  (thợ giỏi)", "≤ 3°C  (mục tiêu)"),
    ("Trường hợp xấu nhất", "+15°C khi phải hạ nhiệt", "Có nhánh hạ nhiệt riêng"),
    ("Độ lặp lại", "Mỗi người một kiểu", "Cùng một luật, mọi ca"),
    ("Dữ liệu để soi lại", "Không có", "CSV 5 giây/dòng"),
]
tw = [Inches(4.0), Inches(3.85), Inches(3.78)]
ty = Inches(2.62)
rh = Inches(0.575)
for r, row in enumerate(rows):
    x = MARGIN
    for c_i, cell in enumerate(row):
        y = ty + rh * r
        if r == 0:
            if cell:
                bx = card(s, x, y, tw[c_i], rh, CARD if c_i == 1 else "2A100E",
                          STROKE if c_i == 1 else RED)
            text(s, x + Inches(0.3), y + Inches(0.19), tw[c_i] - Inches(0.6), Inches(0.3),
                 [(cell, 12, F_BOLD, MUTE if c_i == 1 else RED_BR, 1.6)])
        else:
            if c_i == 0:
                text(s, x + Inches(0.05), y + Inches(0.19), tw[c_i] - Inches(0.3), Inches(0.3),
                     [(cell, 14, F_BODY, DIM)])
            else:
                col = MUTE if c_i == 1 else INK
                fnt = F_BODY if c_i == 1 else F_BOLD
                text(s, x + Inches(0.3), y + Inches(0.17), tw[c_i] - Inches(0.6), Inches(0.35),
                     [(cell, 14.5, fnt, col)])
            rule(s, x, y + rh - Inches(0.01), tw[c_i], "23232B", Pt(1))
        x += tw[c_i] + Inches(0.15)
text(s, MARGIN, Inches(6.35), CW, Inches(0.5),
     [("Số liệu tay lấy từ 7 lần preheat có ghi chép của người vận hành. "
       "Cột tự động là mục tiêu thiết kế — đang ở giai đoạn chạy thử có giám sát.",
       12, F_BODY, DIM)], line=1.3)
page_no(s, 10, TOTAL)

# ─────────────────────────────────────────────────────────────────────────────
# 11 · KẾT
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide("Chốt lại và trạng thái phát hành.")
ember(s, Inches(2.9), Inches(0.4), Inches(8.4), RED, 52)
ember(s, Inches(10.6), Inches(5.0), Inches(5.0), EMBER, 22, "EMBER2")
rule(s, MARGIN, Inches(2.06), Inches(1.5), RED, Pt(5), "RULE")
text(s, MARGIN, Inches(1.5), CW, Inches(0.4),
     [("PREHEAT  ·  OTL-06ALS", 12, F_BOLD, RED, 2.4)], name="KICKER")
text(s, MARGIN, Inches(2.5), Inches(11.0), Inches(2.6),
     [("Bật máy. Đi pha cà phê.", 56, F_DISPLAY, INK, -2),
      ("Quay lại là máy đã nóng.", 56, F_DISPLAY, RED_BR, -2)],
     line=1.1, name="HEADING")
rule(s, MARGIN, Inches(5.3), Inches(11.6), STROKE, Pt(1))
foot = [("Trạng thái", "Chạy thử có giám sát"),
        ("Áp dụng cho", "Đầu đốt thường và premix"),
        ("Liên hệ", "O Tesla  ·  otlpro.com@gmail.com")]
for i, (k, v) in enumerate(foot):
    x = MARGIN + i * Inches(3.95)
    text(s, x, Inches(5.6), Inches(3.7), Inches(0.3),
         [(k.upper(), 10.5, F_BOLD, DIM, 1.6)])
    text(s, x, Inches(5.95), Inches(3.7), Inches(0.4),
         [(v, 15, F_BOLD, INK)])

# ─────────────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "OTL-Preheat-Gioi-thieu.pptx")
prs.save(out)
print("Da luu:", out)
print("So slide:", len(prs.slides.__iter__.__self__._sldIdLst))
