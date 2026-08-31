# -*- coding: utf-8 -*-
"""
Bộ slide kỹ thuật VACUUM CONTROL qua Delta MS300 dành cho PLC.

Nguồn nội dung duy nhất:
  docs/guide/guide-vacuum-control-ms300-plc.md

Đầu ra:
  OTL-Vacuum-Control-MS300-PLC-Masterclass.pptx

Thiết kế sử dụng hình vector và chuyển cảnh Morph để giữ độ sắc nét khi trình chiếu.
"""
import math
import os
import random

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from otl_anim import animate
from otl_deck import Deck


# ─────────────────────────────────────────────────────────────────────────────
# HỆ THIẾT KẾ — “ELECTRIC AIRFLOW”
# ─────────────────────────────────────────────────────────────────────────────
PALETTE = {
    "bg": "050812",
    "bg2": "0A1021",
    "card": "0E1730",
    "card2": "111D3A",
    "cardlo": "091224",
    "stroke": "243457",
    "grid": "111C33",
    "cyan": "15E6FF",
    "blue": "4D7CFF",
    "violet": "8B5CFF",
    "lime": "B9FF4A",
    "amber": "FFC857",
    "red": "FF5576",
    "ink": "F7FAFF",
    "mute": "A6B5D1",
    "dim": "627392",
    "accent": "15E6FF",
    "accent2": "4D7CFF",
}

FONTS = {
    "display": "Segoe UI Black",
    "bold": "Segoe UI Semibold",
    "body": "Segoe UI",
    "light": "Segoe UI Light",
    "num": "Bahnschrift SemiBold Condensed",
    "mono": "Consolas",
}

TOTAL = 20
D = Deck(PALETTE, fonts=FONTS, total=TOTAL, margin=0.72)
M, CW, W, H = D.M, D.CW, D.W, D.H


def IN(value):
    return int(Inches(value))


def no_line(shape):
    shape.line.fill.background()
    return shape


def rect(slide, x, y, w, h, fill="card", alpha=None, radius=False, stroke=None, name=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    if radius:
        shape.adjustments[0] = 0.06
    D.solid(shape, fill, alpha)
    if stroke:
        shape.line.color.rgb = RGBColor.from_string(PALETTE.get(stroke, stroke))
        shape.line.width = Pt(1)
    else:
        no_line(shape)
    if name:
        shape.name = name
    return shape


def circle(slide, cx, cy, diameter, fill="cyan", alpha=None, stroke=None, width=1.0, name=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        int(cx - diameter / 2),
        int(cy - diameter / 2),
        diameter,
        diameter,
    )
    if fill:
        D.solid(shape, fill, alpha)
    else:
        shape.fill.background()
    if stroke:
        shape.line.color.rgb = RGBColor.from_string(PALETTE.get(stroke, stroke))
        shape.line.width = Pt(width)
    else:
        no_line(shape)
    if name:
        shape.name = name
    return shape


def pill(slide, x, y, w, text, color="cyan", fill="cardlo", size=10.5, name=None, upper=True):
    shape = rect(slide, x, y, w, IN(0.34), fill, 92, True, color, name)
    D.text(
        slide, x, y + IN(0.04), w, IN(0.23),
        [(text.upper() if upper else text, size, "bold", color, 1.25)],
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return shape


def micro_grid(slide):
    """Lưới kỹ thuật rất nhẹ để nền có chiều sâu nhưng không gây nhiễu."""
    for i in range(1, 13):
        x = int(W * i / 13)
        D.polyline(slide, [(x, IN(0.38)), (x, H - IN(0.42))], "grid", Pt(0.35)).name = "CHROME_GRID"
    for j in range(1, 7):
        y = int(H * j / 7)
        D.polyline(slide, [(IN(0.35), y), (W - IN(0.35), y)], "grid", Pt(0.35)).name = "CHROME_GRID"


SECTIONS = {
    1: "OPENING",
    2: "WHY",
    3: "WHY",
    4: "SIGNAL",
    5: "SIGNAL",
    6: "SIGNAL",
    7: "SIGNAL",
    8: "SIGNAL",
    9: "CONTROL",
    10: "CONTROL",
    11: "CONTROL",
    12: "LEARNING",
    13: "LEARNING",
    14: "LEARNING",
    15: "OUTPUT",
    16: "PLC",
    17: "PLC",
    18: "SAFETY",
    19: "COMMISSION",
    20: "CLOSING",
}


def _warn_if_title_overflows(n, title, size):
    """Tiêu đề tràn dòng sẽ đè lên phụ đề. Đo bằng font thật thay vì đoán."""
    try:
        from PIL import ImageFont
        f = ImageFont.truetype(r"C:\Windows\Fonts\seguibl.ttf", int(size * 96 / 72))
        w = (f.getlength(title) + (-0.8 * 96 / 72) * max(0, len(title) - 1)) / 96
        pct = w / (W / 914400 - 2 * (M / 914400)) * 100
        if pct > 94:
            print(f"  !! Slide {n}: tieu de chiem {pct:.0f}% be rong -> rut ngan hoac ha title_size")
    except Exception:
        pass


def stage(slide, n, kicker, title, subtitle=None, title_size=38):
    """Khung trình bày lặp lại: section, tiêu đề, thanh tiến độ."""
    _warn_if_title_overflows(n, title, title_size)
    micro_grid(slide)
    D.text(
        slide, M, IN(0.42), IN(5.0), IN(0.23),
        [(f"{SECTIONS[n]}  /  {kicker}".upper(), 10.5, "bold", "cyan", 1.65)],
        name="CHROME_KICKER",
    )
    D.text(
        slide, M, IN(0.78), CW, IN(0.78),
        [(title, title_size, "display", "ink", -0.8)],
        line=0.92,
        name="CHROME_HEADING",
    )
    if subtitle:
        D.text(
            slide, M, IN(1.53), CW, IN(0.42),
            [(subtitle, 13.5, "body", "mute")],
            line=1.2,
            name="CHROME_SUB",
        )
    # Thanh tiến độ chia đoạn tạo cảm giác giao diện thiết bị.
    gap = IN(0.035)
    seg_w = int((CW - gap * (TOTAL - 1)) / TOTAL)
    for i in range(TOTAL):
        x = int(M + i * (seg_w + gap))
        color = "cyan" if i < n else "stroke"
        alpha = 100 if i == n - 1 else (45 if i < n else 65)
        rect(slide, x, H - IN(0.26), seg_w, IN(0.035), color, alpha, name="CHROME_PROG")
    D.text(
        slide, M, H - IN(0.60), IN(4.8), IN(0.22),
        [("O TESLA  ·  VACUUM CONTROL", 9.2, "body", "dim", 0.9)],
        name="CHROME_FOOT",
    )
    D.text(
        slide, W - M - IN(1.70), H - IN(0.60), IN(1.70), IN(0.22),
        [(f"{n:02d}/{TOTAL:02d}", 9.3, "bold", "mute", 0.5)],
        align=PP_ALIGN.RIGHT,
        name="CHROME_FOOT2",
    )


def card_title(slide, x, y, w, tag, title, body, color="cyan", height=IN(1.55)):
    D.card(slide, x, y, w, height, "card", "stroke")
    D.rule(slide, x, y, IN(0.08), color, Pt(5))
    D.text(slide, x + IN(0.28), y + IN(0.15), w - IN(0.52), IN(0.20),
           [(tag.upper(), 9.5, "bold", color, 1.4)])
    D.text(slide, x + IN(0.28), y + IN(0.41), w - IN(0.52), IN(0.31),
           [(title, 16.5, "display", "ink", -0.3)])
    D.text(slide, x + IN(0.28), y + IN(0.79), w - IN(0.52), max(IN(0.31), height - IN(0.82)),
           [(body, 10.7, "body", "mute")], line=1.12)


def node(slide, x, y, w, tag, title, color="cyan", sub=None):
    D.card(slide, x, y, w, IN(1.05), "card", "stroke")
    pill(slide, x + IN(0.18), y + IN(0.15), IN(0.74), tag, color, "cardlo", 9)
    D.text(slide, x + IN(0.18), y + IN(0.56), w - IN(0.36), IN(0.29),
           [(title, 14.2, "bold", "ink")], align=PP_ALIGN.CENTER)
    if sub:
        D.text(slide, x + IN(0.15), y + IN(0.83), w - IN(0.30), IN(0.18),
               [(sub, 9.2, "body", "dim")], align=PP_ALIGN.CENTER)


def flow_chevron(slide, x, y, w=IN(0.34), color="blue", alpha=80):
    ch = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, IN(0.34))
    D.solid(ch, color, alpha)
    no_line(ch)
    return ch


def metric(slide, x, y, w, number, unit, caption, color="cyan"):
    D.text(slide, x, y, w, IN(0.62), [(number, 36, "num", color, -0.7)],
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    D.text(slide, x, y + IN(0.58), w, IN(0.25), [(unit.upper(), 9.5, "bold", "dim", 1.5)],
           align=PP_ALIGN.CENTER)
    D.text(slide, x, y + IN(0.89), w, IN(0.35), [(caption, 11.5, "body", "mute")],
           align=PP_ALIGN.CENTER)


def gauge(slide, cx, cy, diameter, value, max_value, color, label, suffix=""):
    """Đồng hồ vector tối giản; kim chỉ dùng để tạo tương quan thị giác."""
    circle(slide, cx, cy, diameter, "cardlo", 96, "stroke", 1.1)
    circle(slide, cx, cy, int(diameter * 0.76), None, None, color, 3.0)
    angle = math.radians(215 + 250 * max(0.0, min(1.0, value / max_value)))
    r = diameter * 0.31
    px = int(cx + math.cos(angle) * r)
    py = int(cy + math.sin(angle) * r)
    D.polyline(slide, [(cx, cy), (px, py)], color, Pt(3.2))
    circle(slide, cx, cy, IN(0.12), color)
    D.text(slide, cx - diameter // 2, cy - IN(0.28), diameter, IN(0.44),
           [(f"{value:g}{suffix}", 27, "num", "ink", -0.4)], align=PP_ALIGN.CENTER)
    D.text(slide, cx - diameter // 2, cy + IN(0.30), diameter, IN(0.26),
           [(label.upper(), 9, "bold", color, 1.3)], align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 01 · BÌA ĐIỆN ẢNH
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Mở bằng một luận điểm: Air% chỉ là lệnh, Pa mới là kết quả cần giữ.")
micro_grid(s)
D.orb(s, IN(7.0), IN(-1.9), IN(9.6), "blue", 38, "CHROME_FLOW_ORB")
D.orb(s, IN(-2.7), IN(4.55), IN(5.9), "cyan", 18, "CHROME_FLOW_ORB_2")

# Cụm turbine phát sáng ở nửa phải.
cx, cy = IN(10.55), IN(3.46)
circle(s, cx, cy, IN(4.55), "cardlo", 66, "blue", 1.2, "CHROME_TURBINE_RING")
circle(s, cx, cy, IN(3.45), None, None, "cyan", 2.2)
circle(s, cx, cy, IN(2.25), "blue", 12, "violet", 1.0)
for i in range(12):
    ang = math.radians(i * 30 - 90)
    r0, r1 = IN(0.74), IN(1.63)
    p0 = (int(cx + math.cos(ang) * r0), int(cy + math.sin(ang) * r0))
    p1 = (int(cx + math.cos(ang) * r1), int(cy + math.sin(ang) * r1))
    D.polyline(s, [p0, p1], "cyan" if i % 3 == 0 else "blue", Pt(2.0))
circle(s, cx, cy, IN(1.16), "cyan", 92)
D.text(s, cx - IN(1.55), cy - IN(0.145), IN(3.1), IN(0.30),
       [("MS300", 15, "num", "bg", 1.0)], align=PP_ALIGN.CENTER)

D.text(s, M, IN(0.55), IN(6.4), IN(0.28),
       [("OTL-06ALS  /  ENGINEERING MASTERCLASS", 10.5, "bold", "cyan", 1.7)])
D.rule(s, M, IN(1.0), IN(1.45), "cyan", Pt(5), "TITLE_RULE")
D.text(s, M, IN(1.35), IN(7.9), IN(2.4),
       [("VACUUM", 66, "display", "ink", -2.2),
        ("CONTROL", 66, "display", "cyan", -2.2)], line=0.86, name="COVER_HEADING")
D.text(s, M, IN(3.92), IN(7.25), IN(0.83),
       [("Từ 4–20 mA đến một vòng hút tự học trên PLC", 23, "light", "mute")], line=1.15)
D.text(s, M, IN(5.02), IN(6.65), IN(0.55),
       [("Delta MS300  ·  Modbus RTU  ·  S7-1200 / S7-1500", 13, "body", "dim", 0.7)])
pill(s, M, IN(5.75), IN(1.72), "MODBUS RTU", "blue")
pill(s, M + IN(1.90), IN(5.75), IN(1.86), "STEP CONTROL", "cyan")
pill(s, M + IN(3.94), IN(5.75), IN(1.92), "SELF-LEARNING", "lime")
D.text(s, M, H - IN(0.56), IN(5.0), IN(0.24),
       [("SOURCE OF TRUTH  ·  FIRMWARE OTL-06ALS", 9.5, "body", "dim", 1.2)])


# ─────────────────────────────────────────────────────────────────────────────
# 02 · VẤN ĐỀ THẬT
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Hai đồng hồ cùng 40% nhưng lực hút khác nhau; số liệu chỉ là minh hoạ.")
stage(s, 2, "The real problem", "40% gió không đồng nghĩa 40% lực hút", "Cùng một lệnh biến tần, hai trạng thái máy khác nhau cho hai kết quả hoàn toàn khác nhau.")
D.orb(s, IN(9.9), IN(-2.5), IN(6.4), "violet", 20, "CHROME_FLOW_ORB")

left_x, right_x = IN(0.82), IN(6.86)
for x, tag, state, pa, color, notes in [
    (left_x, "MÁY SẠCH", "40% AIR", 146, "cyan", "Lọc sạch · ống thoáng · quạt mát"),
    (right_x, "SAU NHIỀU MẺ", "40% AIR", 92, "red", "Lọc bẩn · bụi bám · tổn thất tăng"),
]:
    D.card(s, x, IN(2.16), IN(5.63), IN(3.75), "card", "stroke")
    pill(s, x + IN(0.35), IN(2.48), IN(1.55), tag, color)
    D.text(s, x + IN(0.35), IN(2.98), IN(2.2), IN(0.56),
           [(state, 33, "num", "ink", -0.8)])
    gauge(s, x + IN(4.18), IN(3.69), IN(2.16), pa, 180, color, "Áp hút", " Pa")
    D.rule(s, x + IN(0.35), IN(4.08), IN(2.05), color, Pt(3))
    D.text(s, x + IN(0.35), IN(4.38), IN(2.65), IN(0.46),
           [(notes, 12.5, "body", "mute")], line=1.2)
    D.text(s, x + IN(0.35), IN(5.13), IN(4.75), IN(0.35),
           [("VÍ DỤ MINH HOẠ", 9.5, "bold", "dim", 1.4)])

D.text(s, M, IN(6.12), CW, IN(0.42),
       [("AIR% là lệnh.  Pa mới là kết quả vật lý tác động lên mẻ rang.", 18, "bold", "lime")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 03 · NORTH STAR
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Đổi mục tiêu điều khiển: giữ áp hút thay vì giữ phần trăm gió.")
stage(s, 3, "Control the outcome", "Đừng giữ lệnh. Hãy giữ kết quả.", "Người vận hành đặt áp hút mục tiêu; hệ thống tự tìm mức gió cần thiết theo trạng thái thật của máy.")
D.orb(s, IN(4.9), IN(1.2), IN(7.7), "blue", 30, "CHROME_FLOW_ORB")

# Trục chuyển đổi lớn.
D.text(s, IN(0.84), IN(2.20), IN(3.4), IN(1.30),
       [("120", 74, "num", "ink", -2.2)], align=PP_ALIGN.CENTER)
D.text(s, IN(0.84), IN(3.30), IN(3.4), IN(0.45),
       [("PA  ·  SETPOINT", 12, "bold", "cyan", 2.0)], align=PP_ALIGN.CENTER)
flow_chevron(s, IN(4.48), IN(2.76), IN(0.72), "cyan")

D.card(s, IN(5.35), IN(2.12), IN(3.10), IN(2.28), "card2", "blue")
D.text(s, IN(5.68), IN(2.46), IN(2.44), IN(0.34),
       [("BỘ NÃO PLC", 12, "bold", "blue", 1.8)], align=PP_ALIGN.CENTER)
D.text(s, IN(5.68), IN(2.91), IN(2.44), IN(0.67),
       [("ĐO  →  LỌC  →  HỌC", 19, "display", "ink", -0.3)], align=PP_ALIGN.CENTER)
D.text(s, IN(5.68), IN(3.60), IN(2.44), IN(0.34),
       [("Và tự bù theo từng mẻ", 11.5, "body", "mute")], align=PP_ALIGN.CENTER)
flow_chevron(s, IN(8.66), IN(2.76), IN(0.72), "lime")

D.text(s, IN(9.55), IN(2.20), IN(2.9), IN(1.15),
       [("43→58", 67, "num", "lime", -2.0)], align=PP_ALIGN.CENTER)
D.text(s, IN(9.55), IN(3.30), IN(2.9), IN(0.45),
       [("% AIR  ·  TỰ THÍCH NGHI", 11.5, "bold", "lime", 1.65)], align=PP_ALIGN.CENTER)

rect(s, M, IN(4.92), CW, IN(0.04), "stroke")
for i, (tag, body, color) in enumerate([
    ("SETPOINT", "Người vận hành chỉ cần nói máy phải giữ bao nhiêu Pa.", "cyan"),
    ("FEEDBACK", "Cảm biến cho biết lực hút thật đang lệch bao nhiêu.", "blue"),
    ("ADAPT", "Air% tự đổi khi lọc bẩn, đường ống đổi hoặc môi trường đổi.", "lime"),
]):
    x = int(M + i * IN(4.03))
    D.text(s, x, IN(5.23), IN(3.75), IN(0.25), [(tag, 10, "bold", color, 1.55)])
    D.text(s, x, IN(5.57), IN(3.75), IN(0.65), [(body, 12.0, "body", "mute")], line=1.24)


# ─────────────────────────────────────────────────────────────────────────────
# 04 · ĐƯỜNG TÍN HIỆU
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Mô tả toàn bộ vòng kín từ buồng rang qua MS300, PLC và trở lại quạt.")
stage(s, 4, "Closed-loop map", "Một vòng kín. Sáu khối. Không có ẩn số.", "MS300 vừa số hoá tín hiệu cảm biến tại cổng ACI, vừa dẫn động quạt theo lệnh đầu ra của PLC.")
D.orb(s, IN(-2.4), IN(4.8), IN(5.5), "cyan", 15, "CHROME_FLOW_ORB")

xs = [IN(0.56), IN(2.68), IN(4.80), IN(6.92), IN(9.04), IN(11.16)]
nodes = [
    ("01", "BUỒNG", "cyan", "Áp âm trong trống"),
    ("02", "CẢM BIẾN", "blue", "4–20 mA"),
    ("03", "MS300 / ACI", "violet", "raw 0…10000"),
    ("04", "MODBUS", "cyan", "FC03 · slave 5"),
    ("05", "PLC", "lime", "PV → Air%"),
    ("06", "QUẠT HÚT", "blue", "cơ cấu chấp hành"),
]
for i, (tag, title, color, sub) in enumerate(nodes):
    node(s, xs[i], IN(2.62), IN(1.60), tag, title, color, sub)
    if i < len(nodes) - 1:
        flow_chevron(s, xs[i] + IN(1.70), IN(2.98), IN(0.29), color)

# Dải tín hiệu nằm dưới từng chặng.
labels = [
    (IN(1.40), "ÁP ÂM", "Pa", "cyan"),
    (IN(3.52), "ANALOG", "4–20 mA", "blue"),
    (IN(5.64), "DIGITAL", "0…10000", "violet"),
    (IN(7.76), "PROCESS", "PV / SP", "cyan"),
    (IN(9.88), "OUTPUT", "Air%", "lime"),
]
for cx0, t, v, color in labels:
    circle(s, cx0, IN(4.58), IN(0.15), color)
    D.text(s, cx0 - IN(0.85), IN(4.20), IN(1.7), IN(0.24),
           [(t, 9.5, "bold", color, 1.35)], align=PP_ALIGN.CENTER)
    D.text(s, cx0 - IN(0.85), IN(4.74), IN(1.7), IN(0.32),
           [(v, 15, "num", "ink", 0.3)], align=PP_ALIGN.CENTER)

D.card(s, IN(2.14), IN(5.44), IN(9.05), IN(0.76), "cardlo", "stroke")
D.text(s, IN(2.45), IN(5.66), IN(8.45), IN(0.30),
       [("VÒNG KHÉP:  quạt → áp hút → cảm biến → PLC → Air% → quạt", 15, "bold", "lime", 0.45)],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 05 · MS300 HAI VAI TRÒ
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Làm rõ vai trò kép của biến tần trong kiến trúc hiện tại.")
stage(s, 5, "One drive · two jobs", "MS300 vừa là giác quan, vừa là cơ bắp", "Cùng một biến tần vừa nhận tín hiệu cảm biến, vừa điều khiển chính chiếc quạt tạo ra áp hút đó.")
D.orb(s, IN(8.9), IN(-2.6), IN(6.7), "violet", 22, "CHROME_FLOW_ORB")

# Khối trung tâm.
D.card(s, IN(4.72), IN(2.10), IN(3.89), IN(3.60), "card2", "blue")
pill(s, IN(5.62), IN(2.43), IN(2.08), "DELTA MS300", "cyan")
circle(s, IN(6.665), IN(3.63), IN(1.55), "blue", 15, "cyan", 2.4)
D.text(s, IN(5.76), IN(3.30), IN(1.81), IN(0.52), [("VFD", 29, "num", "ink", 0.8)], align=PP_ALIGN.CENTER)
D.text(s, IN(5.36), IN(4.66), IN(2.61), IN(0.42), [("SLAVE 5  ·  38400", 12, "mono", "mute")], align=PP_ALIGN.CENTER)

# Vế đo.
D.card(s, IN(0.75), IN(2.34), IN(3.28), IN(3.08), "card", "stroke")
pill(s, IN(1.05), IN(2.65), IN(1.08), "MEASURE", "violet")
D.text(s, IN(1.05), IN(3.20), IN(2.65), IN(0.45), [("ACI  →  ADC", 25, "num", "ink", -0.4)])
D.text(s, IN(1.05), IN(3.82), IN(2.65), IN(0.84),
       [("Nhận 4–20 mA và công bố thành raw 0…10000 qua thanh ghi monitor.", 13.5, "body", "mute")], line=1.25)
pill(s, IN(1.05), IN(4.77), IN(1.46), "0x220C", "cyan", "cardlo", 11)
flow_chevron(s, IN(4.14), IN(3.68), IN(0.42), "violet")

# Vế chấp hành.
D.card(s, IN(9.29), IN(2.34), IN(3.28), IN(3.08), "card", "stroke")
pill(s, IN(9.59), IN(2.65), IN(1.08), "ACTUATE", "lime")
D.text(s, IN(9.59), IN(3.20), IN(2.65), IN(0.45), [("Air%  →  Hz", 25, "num", "ink", -0.4)])
D.text(s, IN(9.59), IN(3.82), IN(2.65), IN(0.84),
       [("Nhận lệnh analog AVI hoặc lệnh tần số Modbus rồi kéo quạt nhanh/chậm.", 13.5, "body", "mute")], line=1.25)
pill(s, IN(9.59), IN(4.77), IN(1.64), "AO / FC06", "lime", "cardlo", 10.5)
flow_chevron(s, IN(8.78), IN(3.68), IN(0.42), "lime")

D.text(s, M, IN(6.07), CW, IN(0.36),
       [("Ưu điểm: đường analog ngắn.  Đánh đổi: mất Modbus là mất luôn giá trị đo.", 15, "bold", "amber")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 06 · ĐỊA CHỈ MODBUS
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Giải thích quy đổi địa chỉ Modbus thô sang địa chỉ 4xxxx của Siemens.")
stage(s, 6, "The address trap", "Một thanh ghi. Ba cách gọi.", "Lệch một đơn vị vẫn trả về con số trông hợp lý — nên bắt buộc kiểm chứng bằng tác động vật lý.")
D.orb(s, IN(9.8), IN(4.3), IN(5.8), "red", 13, "CHROME_FLOW_ORB")

cols = [
    (IN(0.76), "HEX", "0x220C", "MS300 manual / khung RTU", "violet"),
    (IN(4.51), "DECIMAL", "8716", "Firmware ModbusMaster", "cyan"),
    (IN(8.26), "SIEMENS", "48717", "MB_DATA_ADDR 1-based", "lime"),
]
for i, (x, tag, num, cap, color) in enumerate(cols):
    D.card(s, x, IN(2.28), IN(3.32), IN(2.66), "card", color)
    pill(s, x + IN(0.34), IN(2.58), IN(1.12), tag, color)
    D.text(s, x + IN(0.28), IN(3.12), IN(2.76), IN(0.85),
           [(num, 47, "num", "ink", -1.0)], align=PP_ALIGN.CENTER)
    D.text(s, x + IN(0.28), IN(4.18), IN(2.76), IN(0.30),
           [(cap, 11.5, "body", "mute")], align=PP_ALIGN.CENTER)
    if i < 2:
        flow_chevron(s, x + IN(3.43), IN(3.27), IN(0.42), color)

D.card(s, IN(2.42), IN(5.27), IN(8.48), IN(0.78), "cardlo", "stroke")
D.text(s, IN(2.72), IN(5.46), IN(7.88), IN(0.37),
       [("MB_DATA_ADDR  =  40001  +  8716  =  48717", 22, "mono", "cyan", 0.25)],
       align=PP_ALIGN.CENTER)
D.text(s, M, IN(6.23), CW, IN(0.32),
       [("KIỂM CHỨNG: bịt / hở ống áp suất, xác nhận giá trị đổi đúng chiều.", 12.5, "bold", "amber", 1.0)],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 07 · 4–20 mA → RAW → Pa
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Diễn giải chuỗi quy đổi từ dòng điện cảm biến đến đơn vị kỹ thuật Pa.")
stage(s, 7, "Signal normalization", "Từ 4–20 mA thành một con số có ý nghĩa", "MS300 chuẩn hoá tín hiệu về 0…10000; PLC lọc nhiễu rồi nội suy theo minPT / maxPT.")
D.orb(s, IN(-2.7), IN(-1.8), IN(5.5), "blue", 18, "CHROME_FLOW_ORB")

# Dải scale chính.
x0, x1, y = IN(1.16), IN(12.15), IN(3.35)
rect(s, x0, y, x1 - x0, IN(0.12), "stroke")
marks = [
    (0.00, "4 mA", "0", "cyan"),
    (0.25, "8 mA", "2500", "blue"),
    (0.50, "12 mA", "5000", "violet"),
    (1.00, "20 mA", "10000", "lime"),
]
for frac, ma, raw, color in marks:
    x = int(x0 + (x1 - x0) * frac)
    circle(s, x, y + IN(0.06), IN(0.18), color)
    D.text(s, x - IN(0.55), y - IN(0.68), IN(1.1), IN(0.28),
           [(ma, 13, "bold", color, 0.4)], align=PP_ALIGN.CENTER)
    D.text(s, x - IN(0.65), y + IN(0.36), IN(1.3), IN(0.45),
           [(raw, 19, "num", "ink", -0.2)], align=PP_ALIGN.CENTER)

D.text(s, x0, IN(2.14), IN(3.4), IN(0.27), [("DÒNG CẢM BIẾN", 10, "bold", "dim", 1.5)])
D.text(s, x0, IN(4.39), IN(3.4), IN(0.27), [("RAW CỦA MS300", 10, "bold", "dim", 1.5)])

# Công thức nội suy nằm trong khối kính.
D.card(s, IN(2.02), IN(5.07), IN(9.30), IN(0.91), "card2", "cyan")
D.text(s, IN(2.25), IN(5.31), IN(8.84), IN(0.42),
       [("PV  =  minPT  +  raw / 10000  ×  (maxPT − minPT)", 21, "mono", "ink", -0.1)],
       align=PP_ALIGN.CENTER)
pill(s, IN(10.05), IN(5.95), IN(1.28), "OUTPUT: Pa", "lime")

D.text(s, M, IN(6.26), CW, IN(0.26),
       [("BẮT BUỘC tính bằng REAL.  Chia số nguyên sẽ làm PV đứng lì ở minPT.", 12.5, "bold", "red", 0.75)],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 08 · LỌC NHIỄU
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Chứng minh trực quan rằng Kalman một tầng hội tụ thành IIR với alpha xấp xỉ 0,13.")
stage(s, 8, "Noise shaping", "Giữ tín hiệu thật. Loại nhiễu và gai.", "Firmware dùng Kalman vô hướng. Trên PLC, IIR bậc 1 với α = 0,13 cho hành vi xác lập tương đương.")
D.orb(s, IN(9.7), IN(-2.8), IN(6.4), "cyan", 16, "CHROME_FLOW_ORB")

# Đồ thị raw và filtered được tạo từ một chuỗi mô phỏng cố định.
gx0, gx1 = IN(0.78), IN(8.34)
gy0, gy1 = IN(2.28), IN(5.72)
D.card(s, gx0, gy0, gx1 - gx0, gy1 - gy0, "cardlo", "stroke")
px0, px1 = gx0 + IN(0.45), gx1 - IN(0.28)
py0, py1 = gy0 + IN(0.55), gy1 - IN(0.45)
for k in range(5):
    yy = int(py0 + (py1 - py0) * k / 4)
    D.polyline(s, [(px0, yy), (px1, yy)], "grid", Pt(0.6))

random.seed(17)
raw_values = []
for i in range(72):
    target = 2100 if i < 18 else (5200 if i < 49 else 3900)
    noise = random.randint(-520, 520) + int(180 * math.sin(i * 1.41))
    if i in (13, 38, 63):
        noise += 1800 if i != 38 else -1600
    raw_values.append(max(0, min(10000, target + noise)))
filtered_values = []
filt = float(raw_values[0])
for value in raw_values:
    filt += 0.13 * (value - filt)
    filtered_values.append(filt)

lo, hi = 900.0, 7000.0
raw_pts, filt_pts = [], []
for i, value in enumerate(raw_values):
    x = int(px0 + (px1 - px0) * i / (len(raw_values) - 1))
    yr = int(py1 - (py1 - py0) * (value - lo) / (hi - lo))
    yf = int(py1 - (py1 - py0) * (filtered_values[i] - lo) / (hi - lo))
    raw_pts.append((x, yr))
    filt_pts.append((x, yf))
D.polyline(s, raw_pts, "dim", Pt(1.1))
D.polyline(s, filt_pts, "cyan", Pt(3.1))
D.text(s, gx0 + IN(0.28), gy0 + IN(0.16), IN(3.5), IN(0.24),
       [("MÔ PHỎNG CHU KỲ 100 ms", 9.5, "bold", "dim", 1.35)])
pill(s, gx1 - IN(2.62), gy0 + IN(0.16), IN(0.98), "RAW", "dim", "cardlo", 9)
pill(s, gx1 - IN(1.50), gy0 + IN(0.16), IN(1.13), "IIR 0.13", "cyan", "cardlo", 9)

# Khối chứng minh và thông số bên phải.
D.card(s, IN(8.72), IN(2.28), IN(3.82), IN(3.44), "card", "stroke")
pill(s, IN(9.04), IN(2.58), IN(1.68), "KALMAN → IIR", "violet")
D.text(s, IN(9.02), IN(3.18), IN(3.22), IN(0.54),
       [("K∞ ≈ 0,132", 29, "num", "cyan", -0.4)], align=PP_ALIGN.CENTER)
D.text(s, IN(9.02), IN(3.86), IN(3.22), IN(0.68),
       [("filtered := filtered +\n0,13 × (raw − filtered)", 14.5, "mono", "ink")],
       align=PP_ALIGN.CENTER, line=1.15)
D.rule(s, IN(9.20), IN(4.78), IN(2.86), "stroke", Pt(1))
metric(s, IN(9.04), IN(4.91), IN(1.42), "0,77", "giây", "τ ở 100 ms", "lime")
metric(s, IN(10.66), IN(4.91), IN(1.42), "1", "tầng", "không lọc chồng", "amber")

D.text(s, M, IN(6.13), CW, IN(0.37),
       [("Nên đặt median-of-3 trước IIR để loại gai đơn — rẻ, minh bạch và dễ bảo trì.", 13, "bold", "lime")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 09 · TẠI SAO KHÔNG PID
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Giải thích tại sao đối tượng quạt và đường ống phù hợp với feed-forward cộng step controller.")
stage(s, 9, "Control philosophy", "Không dùng PID. Đó là chủ ý thiết kế.", "Feed-forward đảm nhiệm phần thô; step controller đảm nhiệm phần tinh bằng các bước gió nguyên 1%.")
D.orb(s, IN(-2.6), IN(4.6), IN(6.0), "red", 14, "CHROME_FLOW_ORB")

# Cụm PID bị loại bỏ.
D.card(s, IN(0.76), IN(2.24), IN(4.10), IN(3.72), "cardlo", "red")
D.text(s, IN(1.05), IN(2.55), IN(3.52), IN(1.15), [("PID", 66, "num", "dim", -1.6)], align=PP_ALIGN.CENTER)
D.polyline(s, [(IN(1.26), IN(3.51)), (IN(4.34), IN(2.53))], "red", Pt(7))
reasons = [
    ("TRỄ", "Đường khí đáp ứng chậm."),
    ("PHI TUYẾN", "Gain đổi theo vùng Air%."),
    ("NHIỄU", "D khuếch đại xoáy khí."),
    ("1%", "I dễ tích luỹ windup."),
]
for i, (tag, body) in enumerate(reasons):
    x = IN(1.05) + (i % 2) * IN(1.78)
    y = IN(4.02) + (i // 2) * IN(0.80)
    D.text(s, x, y, IN(1.58), IN(0.22), [(tag, 9.5, "bold", "red", 1.0)])
    D.text(s, x, y + IN(0.24), IN(1.58), IN(0.38), [(body, 10.5, "body", "mute")], line=1.10)

# Kiến trúc thay thế ở bên phải.
D.card(s, IN(5.28), IN(2.24), IN(7.28), IN(3.72), "card", "stroke")
pill(s, IN(5.62), IN(2.55), IN(1.94), "TWO-LAYER CONTROL", "lime")

D.card(s, IN(5.64), IN(3.18), IN(2.80), IN(1.74), "card2", "blue")
D.text(s, IN(5.90), IN(3.50), IN(2.28), IN(0.38), [("FEED-FORWARD", 17, "display", "blue", -0.25)], align=PP_ALIGN.CENTER)
D.text(s, IN(5.90), IN(4.02), IN(2.28), IN(0.52), [("Nhảy tới vùng Air% gần đúng đã học.", 12, "body", "mute")], align=PP_ALIGN.CENTER, line=1.2)

flow_chevron(s, IN(8.63), IN(3.86), IN(0.50), "cyan")

D.card(s, IN(9.34), IN(3.18), IN(2.80), IN(1.74), "card2", "cyan")
D.text(s, IN(9.60), IN(3.50), IN(2.28), IN(0.38), [("STEP CONTROL", 17, "display", "cyan", -0.25)], align=PP_ALIGN.CENTER)
D.text(s, IN(9.60), IN(4.02), IN(2.28), IN(0.52), [("Bù nốt sai lệch bằng bước 1% rất êm.", 12, "body", "mute")], align=PP_ALIGN.CENTER, line=1.2)

D.text(s, IN(5.63), IN(5.30), IN(6.55), IN(0.37),
       [("THÔ NHANH  +  TINH CHẮC  =  HỘI TỤ KHÔNG VỌT LỐ", 13, "bold", "lime", 1.0)],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 10 · LUẬT ĐIỀU KHIỂN
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Minh hoạ vùng chết ±3 Pa và hai hướng nhích Air%.")
stage(s, 10, "The control law", "Ba vùng. Một quyết định mỗi lần.", "Chỉ sai lệch ngoài ±3 Pa mới sinh ra bước; mỗi bước thay đổi đúng 1% gió.")
D.orb(s, IN(10.4), IN(4.4), IN(5.8), "cyan", 17, "CHROME_FLOW_ORB")

# Thang sai lệch ngang.
x0, x1, y = IN(1.05), IN(12.28), IN(3.52)
rect(s, x0, y - IN(0.09), x1 - x0, IN(0.18), "stroke")
dead_x0 = int(x0 + (x1 - x0) * 0.43)
dead_x1 = int(x0 + (x1 - x0) * 0.57)
rect(s, dead_x0, y - IN(0.22), dead_x1 - dead_x0, IN(0.44), "lime", 24, True, "lime")
for frac, lab, color in [(0.0, "−30 Pa", "blue"), (0.43, "−3", "lime"), (0.50, "0", "ink"), (0.57, "+3", "lime"), (1.0, "+30 Pa", "cyan")]:
    x = int(x0 + (x1 - x0) * frac)
    circle(s, x, y, IN(0.16), color)
    D.text(s, x - IN(0.55), y + IN(0.31), IN(1.10), IN(0.30), [(lab, 11.5, "num", color)], align=PP_ALIGN.CENTER)

# Ba vùng hành động.
zones = [
    (IN(0.82), "PV > SP", "HÚT QUÁ MẠNH", "Air% −= 1", "blue"),
    (IN(4.73), "|ERROR| ≤ 3", "VÙNG CHẾT", "ĐỨNG YÊN", "lime"),
    (IN(8.64), "PV < SP", "HÚT CHƯA ĐỦ", "Air% += 1", "cyan"),
]
for x, tag, title, action, color in zones:
    D.card(s, x, IN(4.32), IN(3.88), IN(1.42), "card", color)
    D.text(s, x + IN(0.24), IN(4.54), IN(3.40), IN(0.23), [(tag, 9.5, "bold", color, 1.25)], align=PP_ALIGN.CENTER)
    D.text(s, x + IN(0.24), IN(4.83), IN(3.40), IN(0.31), [(title, 15, "display", "ink", -0.2)], align=PP_ALIGN.CENTER)
    D.text(s, x + IN(0.24), IN(5.24), IN(3.40), IN(0.31), [(action, 15.5, "num", color, 0.25)], align=PP_ALIGN.CENTER)

D.text(s, M, IN(6.12), CW, IN(0.35),
       [("error = SP − PV  ·  quy ước bắt buộc: hút mạnh hơn → số Pa lớn hơn", 12.5, "mono", "amber")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 11 · COOLDOWN ĐỘNG
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Biểu diễn cooldown giảm tuyến tính từ 3000 ms xuống 1500 ms theo sai lệch.")
stage(s, 11, "Dynamic pacing", "Xa đi nhanh. Gần đi chậm.", "Biên độ bước luôn là 1%; tốc độ hội tụ điều khiển bằng khoảng nghỉ giữa hai bước.")
D.orb(s, IN(-2.3), IN(-2.4), IN(5.8), "violet", 16, "CHROME_FLOW_ORB")

# Đồ thị cooldown.
gx0, gx1 = IN(0.78), IN(8.35)
gy0, gy1 = IN(2.30), IN(5.74)
D.card(s, gx0, gy0, gx1 - gx0, gy1 - gy0, "cardlo", "stroke")
px0, px1 = gx0 + IN(0.70), gx1 - IN(0.38)
py0, py1 = gy0 + IN(0.58), gy1 - IN(0.55)
for k in range(4):
    yy = int(py0 + (py1 - py0) * k / 3)
    D.polyline(s, [(px0, yy), (px1, yy)], "grid", Pt(0.65))
line_pts = []
for err in range(3, 46):
    clipped = min(err, 30)
    cooldown = 3000 - (clipped - 3) / 27.0 * 1500
    x = int(px0 + (px1 - px0) * (err - 3) / 42.0)
    yv = int(py1 - (py1 - py0) * (cooldown - 1400) / 1700.0)
    line_pts.append((x, yv))
D.polyline(s, line_pts, "cyan", Pt(3.4))
for err, cooldown, color in [(3, 3000, "lime"), (10, 2611, "blue"), (20, 2056, "violet"), (30, 1500, "cyan")]:
    x = int(px0 + (px1 - px0) * (err - 3) / 42.0)
    yv = int(py1 - (py1 - py0) * (cooldown - 1400) / 1700.0)
    circle(s, x, yv, IN(0.13), color)
D.text(s, gx0 + IN(0.28), gy0 + IN(0.16), IN(3.4), IN(0.24), [("COOLDOWN (ms)", 9.5, "bold", "cyan", 1.4)])
D.text(s, gx0 + IN(0.28), gy1 - IN(0.30), gx1 - gx0 - IN(0.56), IN(0.22), [("|ERROR| (Pa)  →", 9.5, "body", "dim", 1.1)], align=PP_ALIGN.CENTER)

# Công thức và nhịp bên phải.
D.card(s, IN(8.72), IN(2.30), IN(3.82), IN(3.44), "card", "stroke")
pill(s, IN(9.04), IN(2.59), IN(1.64), "PACE, NOT GAIN", "violet")
metric(s, IN(8.97), IN(3.14), IN(1.48), "3,0", "giây", "vừa ra khỏi band", "lime")
metric(s, IN(10.70), IN(3.14), IN(1.48), "1,5", "giây", "khi lệch ≥ 30 Pa", "cyan")
D.rule(s, IN(9.10), IN(4.56), IN(3.05), "stroke", Pt(1))
D.text(s, IN(9.02), IN(4.83), IN(3.24), IN(0.63),
       [("Vai trò của P được chuyển từ biên độ bước sang tần suất bước.", 13, "body", "mute")],
       align=PP_ALIGN.CENTER, line=1.22)

D.text(s, M, IN(6.13), CW, IN(0.32),
       [("Giới hạn cứng: 1% / 1,5 s.  Muốn đổi nhanh hơn phải dùng SNAP — không hạ cooldown.", 12.5, "bold", "amber")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 12 · BẢNG FEED-FORWARD
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Giải thích bảng feed-forward Pa sang Air% và cách hệ thống tự học khi ổn định.")
stage(s, 12, "The learned map", "Bảng nhớ biến hệ chậm thành hệ có kinh nghiệm", "Mỗi dòng trả lời đúng một câu hỏi: để giữ mức Pa này, máy thường cần bao nhiêu phần trăm gió?", title_size=34)
D.orb(s, IN(9.7), IN(-2.6), IN(6.3), "blue", 18, "CHROME_FLOW_ORB")

# Đồ thị đặc tuyến Pa → Air%.
gx0, gx1 = IN(0.78), IN(7.76)
gy0, gy1 = IN(2.25), IN(5.84)
D.card(s, gx0, gy0, gx1 - gx0, gy1 - gy0, "cardlo", "stroke")
px0, px1 = gx0 + IN(0.68), gx1 - IN(0.38)
py0, py1 = gy0 + IN(0.56), gy1 - IN(0.55)
for k in range(5):
    yy = int(py0 + (py1 - py0) * k / 4)
    D.polyline(s, [(px0, yy), (px1, yy)], "grid", Pt(0.6))
curve = []
samples = []
for air in range(0, 101, 4):
    pa = 2.0 + 240.0 * (air / 100.0) ** 1.52
    x = int(px0 + (px1 - px0) * air / 100.0)
    yy = int(py1 - (py1 - py0) * pa / 250.0)
    curve.append((x, yy))
    if air in (40, 60, 80):
        samples.append((x, yy, int(pa), air))
D.polyline(s, curve, "cyan", Pt(3.3))
for x, yy, pa, air in samples:
    circle(s, x, yy, IN(0.14), "lime")
    pill(s, x - IN(0.52), yy - IN(0.54), IN(1.04), f"{pa} Pa / {air}%", "lime", "cardlo", 8.2)
D.text(s, gx0 + IN(0.28), gy0 + IN(0.16), IN(3.0), IN(0.22), [("ĐẶC TUYẾN TỰ HỌC", 9.5, "bold", "cyan", 1.35)])
D.text(s, gx0 + IN(0.18), gy1 - IN(0.27), gx1 - gx0 - IN(0.36), IN(0.21), [("AIR%  →", 9.2, "body", "dim", 1.2)], align=PP_ALIGN.CENTER)

# Khối logic học.
D.card(s, IN(8.12), IN(2.25), IN(4.43), IN(3.59), "card", "stroke")
pill(s, IN(8.45), IN(2.56), IN(1.52), "60 ROWS", "blue")
learn = [
    ("10 s", "Ổn định trong ±3 Pa", "cyan"),
    ("≤ 3 Pa", "Tìm dòng gần nhất", "blue"),
    ("≥ 3%", "Drift lớn → học nhanh", "amber"),
    ("60 s", "Có đổi mới lưu xuống nhớ", "lime"),
]
for i, (num, body, color) in enumerate(learn):
    y = IN(3.12) + i * IN(0.63)
    D.text(s, IN(8.45), y, IN(0.90), IN(0.28), [(num, 15, "num", color)], align=PP_ALIGN.RIGHT)
    D.rule(s, IN(9.53), y + IN(0.12), IN(0.28), color, Pt(2.5))
    D.text(s, IN(9.98), y, IN(2.16), IN(0.31), [(body, 12.5, "body", "mute")])

D.text(s, M, IN(6.15), CW, IN(0.31),
       [("Trên PLC: đặt bảng trong DB retentive hoặc recipe — cắt điện không được mất trí nhớ.", 12.5, "bold", "lime")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 13 · SNAP
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Giải thích snap có đệm, ngưỡng 30 Pa và giới hạn nhảy tối đa 20%.")
stage(s, 13, "Fast setpoint changes", "SNAP: nhảy nhanh và cố ý nhảy hụt", "Đổi setpoint lớn thì tra thẳng bảng, trừ hao một vùng đệm, phần còn lại để step controller bù nốt.")
D.orb(s, IN(-2.5), IN(4.2), IN(6.1), "violet", 17, "CHROME_FLOW_ORB")

# Storyboard đổi setpoint.
D.card(s, IN(0.74), IN(2.24), IN(12.02), IN(2.78), "cardlo", "stroke")
steps = [
    (IN(1.08), "SP CŨ", "120 Pa", "blue"),
    (IN(4.02), "SP MỚI", "200 Pa", "cyan"),
    (IN(6.96), "FF TARGET", "68% Air", "violet"),
    (IN(9.90), "SNAP THỰC", "53% Air", "lime"),
]
for i, (x, tag, value, color) in enumerate(steps):
    pill(s, x, IN(2.62), IN(1.46), tag, color)
    D.text(s, x - IN(0.10), IN(3.22), IN(1.68), IN(0.74), [(value, 31, "num", "ink", -0.6)], align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        flow_chevron(s, x + IN(1.78), IN(3.37), IN(0.56), color)
    if i == 2:
        D.text(s, x - IN(0.12), IN(4.16), IN(1.74), IN(0.30), [("− buffer 15%", 11, "bold", "amber")], align=PP_ALIGN.CENTER)

# Hai luật nằm dưới.
card_title(s, IN(0.76), IN(5.20), IN(3.78), "RULE 01", "Chỉ snap khi ΔSP > 30 Pa", "Đổi nhỏ thì giữ Air% hiện tại và bò bằng bước 1%.", "cyan", IN(1.22))
card_title(s, IN(4.78), IN(5.20), IN(3.78), "RULE 02", "Không nhảy quá 20%", "Giới hạn sốc cơ khí cho quạt và đường ống.", "violet", IN(1.22))
card_title(s, IN(8.80), IN(5.20), IN(3.78), "RULE 03", "Hướng theo SP cũ → mới", "Không lấy Air% hiện tại để quyết định hướng snap.", "lime", IN(1.22))


# ─────────────────────────────────────────────────────────────────────────────
# 14 · FACTORY AUTO-TUNE
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Trình bày đúng thời lượng auto-tune: 15 giây cộng 51 điểm nhân 3 giây bằng 2 phút 48 giây.")
stage(s, 14, "Factory auto-tune", "2 phút 48 giây để học cả chiếc máy", "Một lần quét 0→100% dựng 51 điểm đặc tuyến; bảng 60 chỗ nên còn dư biên an toàn.")
D.orb(s, IN(10.0), IN(-2.7), IN(6.2), "lime", 13, "CHROME_FLOW_ORB")

# Timeline thời gian.
D.card(s, IN(0.76), IN(2.22), IN(3.10), IN(3.62), "card", "stroke")
pill(s, IN(1.06), IN(2.53), IN(1.22), "TOTAL", "lime")
D.text(s, IN(1.00), IN(3.10), IN(2.62), IN(0.92), [("02:48", 55, "num", "ink", -1.3)], align=PP_ALIGN.CENTER)
D.text(s, IN(1.00), IN(4.03), IN(2.62), IN(0.29), [("PHÚT : GIÂY", 10, "bold", "lime", 1.6)], align=PP_ALIGN.CENTER)
D.rule(s, IN(1.16), IN(4.55), IN(2.30), "stroke", Pt(1))
D.text(s, IN(1.06), IN(4.82), IN(2.50), IN(0.70),
       [("15 s warm-up\n+ 51 × 3 s đo", 15, "mono", "mute")], align=PP_ALIGN.CENTER, line=1.2)

# Dải cột quét 51 điểm.
bx0, bx1 = IN(4.18), IN(12.55)
baseline = IN(5.27)
height = IN(2.50)
count = 51
bar_gap = IN(0.02)
bar_w = int((bx1 - bx0 - bar_gap * (count - 1)) / count)
for i in range(count):
    frac = i / (count - 1)
    h = int(height * (0.07 + 0.93 * frac ** 1.45))
    x = int(bx0 + i * (bar_w + bar_gap))
    color = "cyan" if i % 5 else "lime"
    rect(s, x, baseline - h, max(IN(0.035), bar_w), h, color, 45 + int(45 * frac))
D.polyline(s, [(bx0, baseline), (bx1, baseline)], "stroke", Pt(1.3))
D.text(s, bx0, IN(2.36), IN(4.0), IN(0.27), [("51 ĐIỂM  ·  BƯỚC 2%", 10, "bold", "cyan", 1.4)])
D.text(s, bx0, baseline + IN(0.22), bx1 - bx0, IN(0.25), [("0% AIR                                      100% AIR", 9.5, "num", "dim", 0.6)], align=PP_ALIGN.CENTER)

# Các pha của mỗi điểm.
phases = [
    ("0 s", "Đặt Air% · xoá tích luỹ", "blue"),
    ("1 s", "Chờ ổn định · không lấy mẫu", "violet"),
    ("2–3 s", "Lấy mẫu Pa · tính trung bình", "cyan"),
]
for i, (tm, body, color) in enumerate(phases):
    x = IN(4.18) + i * IN(2.75)
    D.text(s, x, IN(5.75), IN(0.74), IN(0.27), [(tm, 14, "num", color)])
    D.text(s, x + IN(0.80), IN(5.76), IN(1.78), IN(0.42), [(body, 10.8, "body", "mute")], line=1.15)

D.text(s, M, IN(6.30), CW, IN(0.25),
       [("KHÔNG hạ bước quét xuống 1%: 101 điểm sẽ tràn bảng 60 dòng.", 12.5, "bold", "red", 0.8)],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 15 · CHỌN ĐƯỜNG RA
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("So sánh điều khiển quạt bằng analog và bằng lệnh tần số Modbus.")
stage(s, 15, "Output strategy", "Một Air%. Hai đường ra.", "Cả hai phương án đều khả thi; căn cứ lựa chọn là hành vi của hệ khi RS485 hoặc PLC gặp lỗi.")
D.orb(s, IN(-2.4), IN(-2.5), IN(6.0), "blue", 16, "CHROME_FLOW_ORB")

half = IN(5.78)
# Phương án analog.
D.card(s, IN(0.76), IN(2.24), half, IN(3.72), "card2", "lime")
pill(s, IN(1.08), IN(2.54), IN(1.72), "OPTION A · AO", "lime")
D.text(s, IN(1.08), IN(3.11), IN(5.14), IN(0.42), [("0…100%  →  0…27648  →  0–10 V", 19, "mono", "ink", -0.25)])
benefits_a = [
    ("✓", "Quạt không phụ thuộc nhịp Modbus", "lime"),
    ("✓", "RS485 đứt vẫn giữ được đường lái riêng", "lime"),
    ("!", "Tốn một kênh AO và có sai số analog", "amber"),
]
for i, (mark, text_value, color) in enumerate(benefits_a):
    y = IN(3.83) + i * IN(0.50)
    circle(s, IN(1.28), y + IN(0.10), IN(0.24), color)
    D.text(s, IN(1.20), y - IN(0.03), IN(0.16), IN(0.20), [(mark, 10, "bold", "bg")], align=PP_ALIGN.CENTER)
    D.text(s, IN(1.58), y, IN(4.52), IN(0.27), [(text_value, 12.5, "body", "mute")])
pill(s, IN(1.08), IN(5.44), IN(2.28), "ƯU TIÊN FAIL-SAFE", "lime")

# Phương án Modbus.
D.card(s, IN(6.80), IN(2.24), half, IN(3.72), "card", "blue")
pill(s, IN(7.12), IN(2.54), IN(2.06), "OPTION B · FC06", "blue")
D.text(s, IN(7.12), IN(3.11), IN(5.14), IN(0.42), [("Air%  →  Hz×100  →  register 0x2001", 18, "mono", "ink", -0.22)])
benefits_b = [
    ("✓", "Không cần dây analog; độ phân giải 0,01 Hz", "cyan"),
    ("✓", "Đọc lại được tần số thực tế", "cyan"),
    ("!", "Mất RS485 là mất luôn quyền lái quạt", "red"),
]
for i, (mark, text_value, color) in enumerate(benefits_b):
    y = IN(3.83) + i * IN(0.50)
    circle(s, IN(7.32), y + IN(0.10), IN(0.24), color)
    D.text(s, IN(7.24), y - IN(0.03), IN(0.16), IN(0.20), [(mark, 10, "bold", "bg")], align=PP_ALIGN.CENTER)
    D.text(s, IN(7.62), y, IN(4.52), IN(0.27), [(text_value, 12.5, "body", "mute")])
pill(s, IN(7.12), IN(5.44), IN(2.70), "BẮT BUỘC CÀI TIMEOUT", "red")

D.text(s, M, IN(6.18), CW, IN(0.32),
       [("Không áp dải 30–50 Hz của lồng rang cho quạt.  Dải quạt phải đo trên máy thật.", 12.5, "bold", "amber")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 16 · KIẾN TRÚC PHẦN MỀM PLC
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Kiến trúc khối PLC: máy trạng thái Modbus, khối đọc, khối điều khiển, phân quyền và một điểm xuất duy nhất.")
stage(s, 16, "PLC architecture", "Kiến trúc PLC: dây chuyền dữ liệu một chiều", "Một giao dịch Modbus tại một thời điểm, mỗi khối một trách nhiệm, và chỉ một nơi được ghi đầu ra.", title_size=36)
D.orb(s, IN(9.9), IN(-2.8), IN(6.4), "violet", 17, "CHROME_FLOW_ORB")

flow_y = IN(2.58)
blocks = [
    (IN(0.55), IN(2.08), "01", "MB SCHEDULER", "CASE mbStep", "violet"),
    (IN(2.98), IN(2.08), "02", "FB_VacRead", "raw → filter → Pa", "cyan"),
    (IN(5.41), IN(2.08), "03", "FB_VacCtrl", "SP/PV → Air%", "blue"),
    (IN(7.84), IN(2.08), "04", "AirOwner", "priority arbiter", "lime"),
    (IN(10.27), IN(2.08), "05", "OUTPUT", "AO hoặc FC06", "amber"),
]
for i, (x, w, tag, title, sub, color) in enumerate(blocks):
    D.card(s, x, flow_y, w, IN(1.58), "card", color)
    pill(s, x + IN(0.18), flow_y + IN(0.18), IN(0.56), tag, color, "cardlo", 8.5)
    D.text(s, x + IN(0.16), flow_y + IN(0.65), w - IN(0.32), IN(0.30), [(title, 14.5, "bold", "ink")], align=PP_ALIGN.CENTER)
    D.text(s, x + IN(0.16), flow_y + IN(1.05), w - IN(0.32), IN(0.24), [(sub, 10.5, "mono", color)], align=PP_ALIGN.CENTER)
    if i < len(blocks) - 1:
        flow_chevron(s, x + w + IN(0.07), flow_y + IN(0.62), IN(0.24), color)

# Nhánh dữ liệu retentive và giám sát.
D.card(s, IN(1.55), IN(4.74), IN(3.16), IN(1.03), "cardlo", "stroke")
pill(s, IN(1.82), IN(4.93), IN(1.27), "RETENTIVE DB", "blue", "cardlo", 8.5)
D.text(s, IN(1.82), IN(5.36), IN(2.62), IN(0.26), [("ffSp[] · ffAir[] · ffCnt[]", 11, "mono", "mute")], align=PP_ALIGN.CENTER)
D.polyline(s, [(IN(3.14), IN(4.74)), (IN(6.45), IN(4.30))], "blue", Pt(1.4), "dash")

D.card(s, IN(5.08), IN(4.74), IN(3.16), IN(1.03), "cardlo", "stroke")
pill(s, IN(5.35), IN(4.93), IN(1.15), "HMI / DIAG", "cyan", "cardlo", 8.5)
D.text(s, IN(5.35), IN(5.36), IN(2.62), IN(0.26), [("PV · SP · fault · tableFull", 11, "mono", "mute")], align=PP_ALIGN.CENTER)
D.polyline(s, [(IN(6.66), IN(4.74)), (IN(6.66), IN(4.30))], "cyan", Pt(1.4), "dash")

D.card(s, IN(8.61), IN(4.74), IN(3.16), IN(1.03), "cardlo", "stroke")
pill(s, IN(8.88), IN(4.93), IN(1.19), "WATCHDOG", "red", "cardlo", 8.5)
D.text(s, IN(8.88), IN(5.36), IN(2.62), IN(0.26), [("comm · stuck · range · owner", 11, "mono", "mute")], align=PP_ALIGN.CENTER)
D.polyline(s, [(IN(10.19), IN(4.74)), (IN(10.19), IN(4.30))], "red", Pt(1.4), "dash")

D.text(s, M, IN(6.18), CW, IN(0.30),
       [("MB_MASTER phải chờ DONE / ERROR rồi mới chuyển bước.  Không gọi nhiều slave cùng lúc.", 12.5, "bold", "amber")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 17 · PHÂN QUYỀN AIR%
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Mô tả thứ tự ưu tiên Preheat, Vacuum, Roast và Manual trước một điểm ghi duy nhất.")
stage(s, 17, "Output arbitration", "Bốn chế độ. Một đầu ra duy nhất.", "Phân quyền phải tường minh qua AirOwner; để nhiều khối cùng ghi Air% là nguồn lỗi ngẫu nhiên rất khó truy.")
D.orb(s, IN(-2.5), IN(4.6), IN(5.8), "red", 14, "CHROME_FLOW_ORB")

# Các lớp ưu tiên xếp như bậc thang.
levels = [
    (IN(1.02), IN(2.22), IN(10.74), "03", "PREHEAT", "Quyền tuyệt đối trong chu trình sấy lồng", "red"),
    (IN(1.55), IN(3.10), IN(9.68), "02", "VACUUM CONTROL", "Giữ Pa khi được bật và cảm biến hợp lệ", "cyan"),
    (IN(2.08), IN(3.98), IN(8.62), "01", "ROAST PROFILE", "Phát lại Air% hoặc setpoint theo hồ sơ rang", "blue"),
    (IN(2.61), IN(4.86), IN(7.56), "00", "MANUAL / VR", "Mặc định khi không có chế độ tự động giữ quyền", "dim"),
]
for x, y, w, prio, title, body, color in levels:
    D.card(s, x, y, w, IN(0.70), "card", color)
    pill(s, x + IN(0.16), y + IN(0.18), IN(0.50), prio, color, "cardlo", 8)
    D.text(s, x + IN(0.84), y + IN(0.17), IN(2.40), IN(0.28), [(title, 13.5, "bold", "ink", 0.2)])
    D.text(s, x + IN(3.20), y + IN(0.18), w - IN(3.44), IN(0.27), [(body, 11.5, "body", "mute")], align=PP_ALIGN.RIGHT)

# Điểm ghi duy nhất.
D.card(s, IN(4.30), IN(5.90), IN(4.73), IN(0.54), "card2", "lime")
D.text(s, IN(4.52), IN(6.03), IN(4.30), IN(0.24), [("ONE WRITER  →  airflowCommand  →  OUTPUT", 13, "mono", "lime", 0.35)], align=PP_ALIGN.CENTER)
for _, y, _, _, _, _, color in levels:
    D.polyline(s, [(IN(6.665), y + IN(0.70)), (IN(6.665), IN(5.88))], color, Pt(1.15), "dash")


# ─────────────────────────────────────────────────────────────────────────────
# 18 · FAIL-SAFE
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Các lỗi chính và phản ứng fail-safe khuyến nghị cho bản PLC.")
stage(s, 18, "Fault containment", "Mất số đo không được thành mất lực hút", "Nguyên tắc trên PLC: giữ mức gió an toàn, báo lỗi rõ ràng, và không để đầu ra rơi về 0%.")
D.orb(s, IN(10.1), IN(-2.8), IN(6.3), "red", 17, "CHROME_FLOW_ORB")

# Đồng hồ fail-safe lớn.
D.card(s, IN(0.76), IN(2.22), IN(3.32), IN(3.70), "card2", "lime")
pill(s, IN(1.08), IN(2.53), IN(1.58), "SAFE OUTPUT", "lime")
gauge(s, IN(2.42), IN(4.03), IN(2.26), 60, 100, "lime", "Air on fault", "%")
D.text(s, IN(1.08), IN(5.36), IN(2.68), IN(0.34), [("KHÔNG VỀ 0%", 13, "bold", "red", 1.2)], align=PP_ALIGN.CENTER)

# Ma trận lỗi.
faults = [
    ("5×", "MODBUS FAIL", "Giữ PV cuối; đủ ngưỡng thì fault và chuyển gió an toàn.", "red"),
    ("20s", "RAW STUCK", "Air% đổi mà raw đứng im: nghi cảm biến hoặc ống đo lỗi.", "amber"),
    ("0", "RAW AT ZERO", "Quạt >50% mà raw bám 0: nghi đứt dây 4–20 mA.", "violet"),
    ("60", "TABLE FULL", "Bảng FF đầy phải báo HMI; không được âm thầm bỏ điểm.", "blue"),
]
for i, (num, title, body, color) in enumerate(faults):
    x = IN(4.42) + (i % 2) * IN(4.06)
    y = IN(2.22) + (i // 2) * IN(1.85)
    D.card(s, x, y, IN(3.74), IN(1.56), "card", "stroke")
    D.text(s, x + IN(0.25), y + IN(0.20), IN(0.72), IN(0.44), [(num, 27, "num", color, -0.4)], align=PP_ALIGN.CENTER)
    D.text(s, x + IN(1.10), y + IN(0.20), IN(2.31), IN(0.27), [(title, 11, "bold", color, 1.0)])
    D.text(s, x + IN(1.10), y + IN(0.55), IN(2.31), IN(0.68), [(body, 11.5, "body", "mute")], line=1.18)

D.text(s, M, IN(6.18), CW, IN(0.30),
       [("FAIL-SAFE là một trạng thái điều khiển có chủ đích — không phải giá trị mặc định sau lỗi.", 12.5, "bold", "lime")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 19 · COMMISSIONING GATES
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Năm cổng kiểm tra trước khi đưa vòng hút vào vận hành và bốn mục chưa xác nhận phải đóng lại.")
stage(s, 19, "Release gates", "Chỉ bàn giao khi đã vượt đủ năm cổng", "Mỗi cổng phải có bằng chứng đo hoặc thử lỗi. Không chấp nhận kết luận “chắc là đúng”.")
D.orb(s, IN(-2.2), IN(-2.5), IN(5.5), "amber", 14, "CHROME_FLOW_ORB")

gates = [
    ("01", "HARDWARE", "ACI đúng dòng/áp · shield một đầu · 120Ω hai đầu", "cyan"),
    ("02", "DRIVE", "Slave 5 · 38400 · ACI không tranh quyền tần số", "blue"),
    ("03", "DATA", "48717 trả 0…10000 · bịt/hở đổi đúng chiều", "violet"),
    ("04", "TUNE", "51 điểm · bảng đơn điệu · dữ liệu còn sau mất điện", "lime"),
    ("05", "FAULT", "Rút RS485 · rút sensor · xác nhận gió không về 0", "red"),
]
gw = IN(2.34)
for i, (num, title, body, color) in enumerate(gates):
    x = IN(0.62) + i * IN(2.54)
    D.card(s, x, IN(2.18), gw, IN(2.40), "card", color)
    pill(s, x + IN(0.22), IN(2.46), IN(0.53), num, color, "cardlo", 8)
    D.text(s, x + IN(0.22), IN(2.97), gw - IN(0.44), IN(0.31), [(title, 14.5, "bold", "ink", 0.4)], align=PP_ALIGN.CENTER)
    D.rule(s, x + IN(0.55), IN(3.43), gw - IN(1.10), color, Pt(2.6))
    D.text(s, x + IN(0.22), IN(3.73), gw - IN(0.44), IN(0.66), [(body, 10.4, "body", "mute")], align=PP_ALIGN.CENTER, line=1.18)

# Các mục chưa kiểm chứng được giữ như thẻ đỏ trước phát hành.
D.card(s, IN(1.07), IN(4.96), IN(11.18), IN(1.15), "cardlo", "red")
pill(s, IN(1.35), IN(5.19), IN(1.74), "RED TAGS", "red")
red_tags = ["0x220C đúng model", "Khung 8-N-1", "Mã P03-xx của ACI", "Dải Hz thật của quạt"]
for i, label in enumerate(red_tags):
    x = IN(3.38) + i * IN(2.10)
    circle(s, x, IN(5.50), IN(0.16), "red")
    D.text(s, x + IN(0.18), IN(5.37), IN(1.70), IN(0.35), [(label, 10.8, "bold", "mute")], anchor=MSO_ANCHOR.MIDDLE)

D.text(s, M, IN(6.24), CW, IN(0.25),
       [("Căn cứ cuối cùng khi có tranh luận: manual đúng model MS300 và phép đo trên máy thật.", 12, "bold", "amber")],
       align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# 20 · KẾT
# ─────────────────────────────────────────────────────────────────────────────
s = D.slide("Kết thúc bằng ba nhịp nhớ và lời nhắc xác minh trước khi phát hành.")
micro_grid(s)
D.orb(s, IN(6.9), IN(-2.0), IN(9.7), "blue", 39, "CHROME_FLOW_ORB")
D.orb(s, IN(-2.6), IN(4.7), IN(5.8), "cyan", 18, "CHROME_FLOW_ORB_2")

# Lặp lại turbine của bìa để tạo vòng tròn thị giác khi Morph.
cx, cy = IN(10.55), IN(3.46)
circle(s, cx, cy, IN(4.55), "cardlo", 66, "blue", 1.2, "CHROME_TURBINE_RING")
circle(s, cx, cy, IN(3.45), None, None, "cyan", 2.2)
circle(s, cx, cy, IN(2.25), "blue", 12, "violet", 1.0)
for i in range(12):
    ang = math.radians(i * 30 - 90)
    r0, r1 = IN(0.74), IN(1.63)
    p0 = (int(cx + math.cos(ang) * r0), int(cy + math.sin(ang) * r0))
    p1 = (int(cx + math.cos(ang) * r1), int(cy + math.sin(ang) * r1))
    D.polyline(s, [p0, p1], "cyan" if i % 3 == 0 else "blue", Pt(2.0))
circle(s, cx, cy, IN(1.30), "lime", 92)
D.text(s, cx - IN(1.55), cy - IN(0.155), IN(3.1), IN(0.31), [("CLOSED LOOP", 11, "num", "bg", 1.0)], align=PP_ALIGN.CENTER)

D.text(s, M, IN(0.62), IN(6.7), IN(0.28), [("VACUUM CONTROL  /  READY TO IMPLEMENT", 10.5, "bold", "cyan", 1.7)])
D.rule(s, M, IN(1.04), IN(1.45), "cyan", Pt(5), "TITLE_RULE")
D.text(s, M, IN(1.45), IN(7.65), IN(3.14),
       [("Một lần dựng.", 43, "display", "ink", -1.3),
        ("Một lần tune.", 43, "display", "blue", -1.3),
        ("Mọi mẻ đúng lực hút.", 43, "display", "lime", -1.3)],
       line=0.98, name="CLOSE_HEADING")
D.text(s, M, IN(4.78), IN(6.45), IN(0.72),
       [("Air% sẽ thay đổi.  Điều kiện máy sẽ thay đổi.\nPa mục tiêu thì không.", 18, "light", "mute")], line=1.28)
pill(s, M, IN(5.85), IN(1.64), "48717 · FC03", "violet")
pill(s, M + IN(1.84), IN(5.85), IN(1.52), "α = 0,13", "cyan", upper=False)
pill(s, M + IN(3.56), IN(5.85), IN(1.60), "±3 Pa · 1%", "lime")
D.text(s, M, H - IN(0.57), IN(7.0), IN(0.24),
       [("Nguồn: guide-vacuum-control-ms300-plc.md  ·  O Tesla", 9.5, "body", "dim", 1.0)])
D.text(s, W - M - IN(1.25), H - IN(0.57), IN(1.25), IN(0.24),
       [("20 / 20", 10, "num", "dim", 1.1)], align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# XUẤT FILE
# ─────────────────────────────────────────────────────────────────────────────
OUTPUT = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "OTL-Vacuum-Control-MS300-PLC-Masterclass.pptx",
)
D.save(OUTPUT)

# Bơm hiệu ứng động; hình có tên bắt đầu bằng CHROME được bỏ qua
# vì chúng là mỏ neo cho chuyển cảnh Morph.
animate(OUTPUT, skip_prefix="CHROME")
print("Saved:", OUTPUT)
