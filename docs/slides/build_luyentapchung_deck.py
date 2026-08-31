# -*- coding: utf-8 -*-
"""
TOÁN LỚP 3 — LUYỆN TẬP CHUNG (tính giá trị của biểu thức)
Giáo viên: Nguyễn Thị Thanh Tâm · Lớp 3.10 · Trường Tiểu học Tân Bình

Dựng lại từ bản gốc "thi thi toan 13 - 14.ppt" (21 slide, khổ 4:3, năm 2013):
giữ NGUYÊN mạch dạy và toàn bộ số liệu bài tập, làm mới hoàn toàn phần nhìn,
chuyển sang khổ 16:9, bỏ 2 slide trống.

Chạy:  python build_luyentapchung_deck.py
Ra:    F:\\PPT\\Toan3-Luyen-tap-chung-NguyenThiThanhTam.pptx
"""
import math
import os

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from otl_deck import Deck

# ─────────────────────────────────────────────────────────────────────────────
PALETTE = {
    "bg":       "F6FAFF",
    "bg2":      "FFFFFF",
    "card":     "FFFFFF",
    "cardsoft": "EDF4FF",
    "cardwarm": "FFF6E8",
    "cardmint": "E9F8F1",
    "stroke":   "DBE6F6",
    "accent":   "2563EB",   # xanh dương — chủ đạo
    "accent2":  "0E9F6E",   # xanh lá — kết quả đúng
    "coral":    "FF6B4A",   # cam san hô — nhấn
    "gold":     "F5A524",
    "grape":    "7C5CFF",
    "ink":      "132038",
    "mute":     "56668260"[:6] or "566682",
    "dim":      "94A3BE",
    "grid":     "E6EDF8",
    "petal":    "FF7AA2",
}
PALETTE["mute"] = "566682"

D = Deck(PALETTE, total=19)
M, CW, W, H = D.M, D.CW, D.W, D.H
TONG = 19


def IN(v):
    return int(Inches(v))


def rgb(k):
    return RGBColor.from_string(PALETTE[k])


# ─────────────────────────────────────────────────────────────────────────────
# THÀNH PHẦN DÙNG LẠI
# ─────────────────────────────────────────────────────────────────────────────
def foot(s, n):
    D.text(s, M, H - IN(0.68), CW, IN(0.3),
           [("Toán 3  ·  Luyện tập chung  ·  Lớp 3.10", 10.5, "body", "dim")])
    D.text(s, M, H - IN(0.68), CW, IN(0.3),
           [("%d / %d" % (n, TONG), 10.5, "num", "dim", 1.2)], align=PP_ALIGN.RIGHT)


def chip(s, txt, color="accent"):
    """Nhãn bo tròn góc phải trên — cho biết đang ở hoạt động nào."""
    w = IN(0.42) + int(len(txt) * IN(0.105))
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           int(W - M - w), IN(0.66), w, IN(0.44))
    c.adjustments[0] = 0.45
    D.solid(c, color)
    D._no_line(c)
    D.text(s, int(W - M - w), IN(0.78), w, IN(0.3),
           [(txt.upper(), 12, "bold", "bg2", 1.6)], align=PP_ALIGN.CENTER)


def lesson(s, chip_txt=None, chip_col="accent"):
    D.kicker(s, "Toán  ·  Luyện tập chung", color="accent")
    if chip_txt:
        chip(s, chip_txt, chip_col)


def expr(s, left, top, width, lines, size=30, gap=0.5, hi_last=True):
    """Một biểu thức kèm các bước tính. lines[0] là đề, các dòng sau là bước."""
    y = top
    for i, ln in enumerate(lines):
        last = (i == len(lines) - 1)
        col = "accent2" if (last and hi_last and i > 0) else "ink"
        fnt = "display" if (i == 0 or (last and hi_last and i > 0)) else "num"
        D.text(s, left, y, width, IN(gap), [(ln, size, fnt, col, -0.5)])
        y += IN(gap)
    return y


def flower(s, cx, cy, r, petal="petal", core="gold"):
    for k in range(5):
        a = math.radians(90 + k * 72)
        px = cx + math.cos(a) * r * 0.60
        py = cy - math.sin(a) * r * 0.60
        p = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               int(px - r * 0.52), int(py - r * 0.52),
                               int(r * 1.04), int(r * 1.04))
        D.solid(p, petal)
        D._no_line(p)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx - r * 0.42), int(cy - r * 0.42),
                           int(r * 0.84), int(r * 0.84))
    D.solid(c, core)
    D._no_line(c)


def star(s, cx, cy, r, color="gold"):
    st = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                            int(cx - r), int(cy - r), int(2 * r), int(2 * r))
    D.solid(st, color)
    D._no_line(st)


def pill(s, left, top, width, height, txt, size, stroke="stroke", fill="card",
         color="ink"):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    p.adjustments[0] = 0.5
    D.solid(p, fill)
    p.line.color.rgb = rgb(stroke)
    p.line.width = Pt(2)
    D.text(s, left, top + int(height * 0.24), width, IN(0.5),
           [(txt, size, "display", color, -0.5)], align=PP_ALIGN.CENTER)
    return p


# ═════════════════════════════════════════════════════════════════════════════
# 01 · BÌA
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Chào mừng quý thầy cô về dự giờ. Môn Toán, lớp 3.10.")
D.orb(s, IN(8.4), IN(-1.8), IN(8.6), "accent", 22)
D.orb(s, IN(-2.4), IN(4.2), IN(6.0), "gold", 16, "ORB2")
D.text(s, M, IN(1.15), CW, IN(0.4),
       [("TRƯỜNG TIỂU HỌC TÂN BÌNH", 14, "bold", "accent", 3.0)], name="KICKER")
D.rule(s, M, IN(1.72), IN(1.6), "coral", Pt(5), "RULE")
D.text(s, M, IN(2.2), IN(11.6), IN(2.2),
       [("Chào mừng quý thầy cô", 54, "display", "ink", -1.5),
        ("về dự giờ thăm lớp", 54, "display", "accent", -1.5)],
       line=1.12, name="HEADING")
for i, (lab, val, col) in enumerate([("MÔN", "Toán", "coral"), ("LỚP", "3.10", "grape")]):
    x = int(M + i * IN(2.6))
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, IN(4.55), IN(2.3), IN(0.85))
    c.adjustments[0] = 0.28
    D.solid(c, "card")
    c.line.color.rgb = rgb(col)
    c.line.width = Pt(2)
    D.text(s, x + IN(0.3), IN(4.72), IN(0.9), IN(0.3), [(lab, 10.5, "bold", "dim", 1.6)])
    D.text(s, x + IN(0.3), IN(4.95), IN(1.8), IN(0.4), [(val, 22, "display", col, -0.5)])
D.rule(s, M, IN(5.75), IN(11.6), "stroke", Pt(1))
D.text(s, M, IN(6.0), IN(11.6), IN(0.35),
       [("GIÁO VIÊN THỰC HIỆN", 10.5, "bold", "dim", 1.8)])
D.text(s, M, IN(6.3), IN(11.6), IN(0.45),
       [("Nguyễn Thị Thanh Tâm", 22, "display", "ink", -0.5)])
for i in range(3):
    flower(s, IN(10.6 + i * 0.95), IN(5.3), IN(0.3),
           ["petal", "gold", "coral"][i])
star(s, IN(11.3), IN(2.5), IN(0.18), "gold")
star(s, IN(12.3), IN(3.4), IN(0.13), "coral")

# ═════════════════════════════════════════════════════════════════════════════
# 02 · KHỞI ĐỘNG
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Khởi động: cả lớp cùng hát.")
D.orb(s, IN(3.4), IN(0.6), IN(7.6), "gold", 20)
D.kicker(s, "Khởi động", color="gold")
chip(s, "Khởi động", "gold")
D.text(s, M, IN(2.4), IN(11.6), IN(1.6),
       [("Hãy hát lên bạn ơi !", 62, "display", "ink", -1.5)], name="HEADING")
D.text(s, M, IN(4.1), IN(10.0), IN(0.5),
       [("Cả lớp cùng hát một bài thật vui để bắt đầu tiết học nhé.",
         24, "light", "mute")])
for i, (cx, r, col) in enumerate([(5.6, 0.34, "coral"), (7.0, 0.26, "gold"),
                                  (8.4, 0.30, "petal"), (9.8, 0.24, "grape")]):
    flower(s, IN(cx), IN(5.5), IN(r), col)
foot(s, 2)

# ═════════════════════════════════════════════════════════════════════════════
# 03 · KIỂM TRA BÀI CŨ — ĐỀ
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Kiểm tra bài cũ: tính giá trị của hai biểu thức.")
D.orb(s, IN(10.0), IN(-2.0), IN(6.4), "grape", 14)
D.kicker(s, "Kiểm tra bài cũ", color="grape")
chip(s, "Bài cũ", "grape")
D.heading(s, "Tính giá trị của biểu thức")
half = (CW - IN(0.4)) / 2
for i, (lab, e) in enumerate([("a)", "67 – 27 + 10"), ("b)", "9 + 90 : 9")]):
    x = int(M + i * (half + IN(0.4)))
    D.card(s, x, IN(2.85), int(half), IN(2.5), "card", "stroke")
    D.text(s, x + IN(0.55), IN(3.15), IN(0.6), IN(0.4), [(lab, 20, "bold", "grape")])
    D.text(s, x + IN(0.55), IN(3.75), int(half - IN(1.1)), IN(0.8),
           [(e, 42, "display", "ink", -1)])
foot(s, 3)

# ═════════════════════════════════════════════════════════════════════════════
# 04 · KIỂM TRA BÀI CŨ — ĐÁP ÁN
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Đáp án kiểm tra bài cũ.")
D.orb(s, IN(10.0), IN(-2.0), IN(6.4), "grape", 14)
D.kicker(s, "Kiểm tra bài cũ", color="grape")
chip(s, "Đáp án", "accent2")
D.heading(s, "Tính giá trị của biểu thức")
sols = [("a)", ["67 – 27 + 10", "=  40 + 10", "=  50"]),
        ("b)", ["9 + 90 : 9", "=  9 + 10", "=  19"])]
for i, (lab, lines) in enumerate(sols):
    x = int(M + i * (half + IN(0.4)))
    D.card(s, x, IN(2.85), int(half), IN(2.9), "card", "accent2")
    D.text(s, x + IN(0.55), IN(3.1), IN(0.6), IN(0.4), [(lab, 20, "bold", "grape")])
    expr(s, x + IN(0.55), IN(3.6), int(half - IN(1.1)), lines, 34, 0.62)
foot(s, 4)

# ═════════════════════════════════════════════════════════════════════════════
# 05 · BÀI MỚI
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Giới thiệu bài mới và năm nhiệm vụ của tiết học.")
D.orb(s, IN(2.8), IN(-2.2), IN(7.6), "accent", 18)
D.text(s, M, IN(1.0), CW, IN(0.35),
       [("Thứ  ……… ,  ngày  ……  tháng  ……  năm  ………", 13, "body", "dim", 1)],
       name="KICKER")
D.rule(s, M, IN(1.5), IN(1.6), "coral", Pt(5), "RULE")
D.text(s, M, IN(1.85), IN(11.6), IN(0.6),
       [("TOÁN", 20, "display", "coral", 4.0)])
D.text(s, M, IN(2.45), IN(11.6), IN(1.2),
       [("Luyện tập chung", 54, "display", "ink", -1.5)], name="HEADING")
tasks = [("Bài 1", "Tính giá trị\nbiểu thức", "accent"),
         ("Bài 2", "Nhân chia trước\ncộng trừ sau", "coral"),
         ("Bài 3", "Biểu thức\ncó dấu ngoặc", "grape"),
         ("Bài 4", "Nối biểu thức\nvới giá trị", "gold"),
         ("Bài 5", "Giải bài toán\ncó lời văn", "accent2")]
tw = (CW - IN(0.6)) / 5
for i, (t, d, col) in enumerate(tasks):
    x = int(M + i * (tw + IN(0.15)))
    D.card(s, x, IN(4.15), int(tw), IN(1.95))
    D.rule(s, x, IN(4.15), int(tw), col, Pt(4))
    D.text(s, x + IN(0.3), IN(4.45), int(tw - IN(0.6)), IN(0.45),
           [(t, 24, "display", col, -0.5)])
    D.text(s, x + IN(0.3), IN(5.0), int(tw - IN(0.6)), IN(0.9),
           [(d.replace("\n", " "), 14, "body", "mute")], line=1.3)
foot(s, 5)

# ═════════════════════════════════════════════════════════════════════════════
# 06–07 · BÀI 1
# ═════════════════════════════════════════════════════════════════════════════
B1 = [("a)", [["324 – 20 + 61", "=  304 + 61", "=  365"],
              ["21 x 3 : 9", "=  63 : 9", "=  7"]]),
      ("b)", [["188 + 12 – 50", "=  200 – 50", "=  150"],
              ["40 : 2 x 6", "=  20 x 6", "=  120"]])]


def bai1(show):
    ss = D.slide("Bài 1: tính giá trị của biểu thức." + (" Đáp án." if show else ""))
    D.orb(ss, IN(10.6), IN(4.4), IN(6.0), "accent", 12)
    lesson(ss, "Đáp án" if show else "Bài 1", "accent2" if show else "accent")
    D.heading(ss, "Bài 1.  Tính giá trị của biểu thức")
    hw = (CW - IN(0.4)) / 2
    for i, (lab, exprs) in enumerate(B1):
        x = int(M + i * (hw + IN(0.4)))
        D.card(ss, x, IN(2.72), int(hw), IN(3.5), "card", "accent2" if show else "stroke")
        D.text(ss, x + IN(0.5), IN(3.05), IN(0.6), IN(0.4),
               [(lab, 19, "bold", "accent")])
        D.rule(ss, x + IN(0.5), IN(4.72), int(hw - IN(1.0)), "stroke", Pt(1))
        for j, lines in enumerate(exprs):
            y = IN(3.28) + int(IN(1.56) * j)
            expr(ss, x + IN(0.5), y, int(hw - IN(1.0)),
                 lines if show else lines[:1], 28, 0.44)
    return ss


foot(bai1(False), 6)
foot(bai1(True), 7)

# ═════════════════════════════════════════════════════════════════════════════
# 08–09 · BÀI 2
# ═════════════════════════════════════════════════════════════════════════════
B2 = [("a)", ["15 + 7 x 8", "=  15 + 56", "=  71"]),
      ("b)", ["90 + 28 : 2", "=  90 + 14", "=  104"])]


def bai2(show):
    ss = D.slide("Bài 2: biểu thức có cả cộng trừ và nhân chia."
                 + (" Đáp án." if show else ""))
    D.orb(ss, IN(-2.4), IN(-1.6), IN(6.4), "coral", 12)
    lesson(ss, "Đáp án" if show else "Bài 2", "accent2" if show else "coral")
    D.heading(ss, "Bài 2.  Tính giá trị của biểu thức")
    hw = (CW - IN(0.4)) / 2
    for i, (lab, lines) in enumerate(B2):
        x = int(M + i * (hw + IN(0.4)))
        D.card(ss, x, IN(2.8), int(hw), IN(2.6), "card", "accent2" if show else "stroke")
        D.text(ss, x + IN(0.55), IN(3.05), IN(0.6), IN(0.4),
               [(lab, 20, "bold", "coral")])
        expr(ss, x + IN(0.55), IN(3.55), int(hw - IN(1.1)),
             lines if show else lines[:1], 34, 0.6)
    if show:
        b = D.card(ss, M, IN(5.65), CW, IN(0.85), "cardwarm", "coral")
        D.text(ss, M + IN(0.5), IN(5.88), int(CW - IN(1.0)), IN(0.4),
               [("Ghi nhớ:  biểu thức có cộng, trừ, nhân, chia thì làm "
                 "NHÂN CHIA trước, CỘNG TRỪ sau.", 18, "bold", "ink")])
    return ss


foot(bai2(False), 8)
foot(bai2(True), 9)

# ═════════════════════════════════════════════════════════════════════════════
# 10–11 · BÀI 3
# ═════════════════════════════════════════════════════════════════════════════
B3 = [("a)", ["123 x (42 – 40)", "=  123 x 2", "=  246"]),
      ("b)", ["72 : (2 x 4)", "=  72 : 8", "=  9"])]


def bai3(show):
    ss = D.slide("Bài 3: biểu thức có dấu ngoặc." + (" Đáp án." if show else ""))
    D.orb(ss, IN(10.2), IN(-2.2), IN(6.2), "grape", 12)
    lesson(ss, "Đáp án" if show else "Làm phiếu", "accent2" if show else "grape")
    D.heading(ss, "Bài 3.  Tính giá trị của biểu thức")
    hw = (CW - IN(0.4)) / 2
    for i, (lab, lines) in enumerate(B3):
        x = int(M + i * (hw + IN(0.4)))
        D.card(ss, x, IN(2.8), int(hw), IN(2.6), "card", "accent2" if show else "stroke")
        D.text(ss, x + IN(0.55), IN(3.05), IN(0.6), IN(0.4),
               [(lab, 20, "bold", "grape")])
        expr(ss, x + IN(0.55), IN(3.55), int(hw - IN(1.1)),
             lines if show else lines[:1], 34, 0.6)
    if show:
        D.card(ss, M, IN(5.65), CW, IN(0.85), "cardsoft", "grape")
        D.text(ss, M + IN(0.5), IN(5.88), int(CW - IN(1.0)), IN(0.4),
               [("Ghi nhớ:  biểu thức có dấu ngoặc ( ) thì làm TRONG NGOẶC trước.",
                 18, "bold", "ink")])
    return ss


foot(bai3(False), 10)
foot(bai3(True), 11)

# ═════════════════════════════════════════════════════════════════════════════
# 12–13 · BÀI 4 — NỐI
# ═════════════════════════════════════════════════════════════════════════════
LEFT = ["86 – (81 – 31)", "90 + 70 x 2", "142 – (42 : 2)",
        "56 x (17 – 12)", "(142 – 42) : 2"]
RIGHT = ["230", "36", "280", "50", "121"]
MAP = {0: 1, 1: 0, 2: 4, 3: 2, 4: 3}   # trái → phải
LX, LW = IN(1.15), IN(4.4)
RX, RW = IN(8.5), IN(2.5)
YS = [IN(2.6) + int(IN(0.8) * i) for i in range(5)]
PH = IN(0.66)


def bai4(show):
    ss = D.slide("Bài 4: nối biểu thức với giá trị của nó."
                 + (" Đáp án." if show else ""))
    D.orb(ss, IN(-2.2), IN(4.6), IN(5.6), "gold", 12)
    lesson(ss, "Đáp án" if show else "Bài 4", "accent2" if show else "gold")
    D.heading(ss, "Bài 4.  Nối biểu thức với giá trị của nó")
    if show:                       # vẽ đường nối trước để nằm dưới các ô
        for a, b in MAP.items():
            D.polyline(ss, [(LX + LW, YS[a] + PH // 2), (RX, YS[b] + PH // 2)],
                       "accent2", Pt(2.5))
    for i, e in enumerate(LEFT):
        pill(ss, LX, YS[i], LW, PH, e, 26,
             "accent2" if show else "stroke")
    for i, v in enumerate(RIGHT):
        pill(ss, RX, YS[i], RW, PH, v, 28,
             "accent2" if show else "gold",
             "cardmint" if show else "cardwarm",
             "accent2" if show else "coral")
    return ss


foot(bai4(False), 12)
foot(bai4(True), 13)

# ═════════════════════════════════════════════════════════════════════════════
# 14 · BÀI 5 — ĐỀ + TÓM TẮT
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Bài 5: bài toán có lời văn về xếp bánh vào hộp và thùng.")
D.orb(s, IN(10.4), IN(4.2), IN(5.8), "accent2", 12)
lesson(s, "Làm vở", "accent2")
D.heading(s, "Bài 5.  Bài toán")
D.card(s, M, IN(2.7), CW, IN(1.35), "cardsoft", "accent")
D.text(s, M + IN(0.5), IN(2.95), int(CW - IN(1.0)), IN(1.0),
       [("Người ta xếp 800 cái bánh vào các hộp, mỗi hộp có 4 cái. Sau đó xếp các hộp "
         "vào thùng, mỗi thùng có 5 hộp. Hỏi có bao nhiêu thùng bánh ?",
         22, "body", "ink")], line=1.35)
D.text(s, M, IN(4.35), CW, IN(0.4), [("TÓM TẮT", 13, "bold", "dim", 2.2)])
tt = [("4 cái bánh", "1 hộp"), ("800 cái bánh", "…… hộp ?"),
      ("5 hộp", "1 thùng"), ("Có", "…… thùng ?")]
sw = (CW - IN(0.45)) / 4
for i, (a, b) in enumerate(tt):
    x = int(M + i * (sw + IN(0.15)))
    D.card(s, x, IN(4.75), int(sw), IN(1.35), "card", "stroke")
    D.text(s, x + IN(0.3), IN(4.98), int(sw - IN(0.6)), IN(0.4),
           [(a, 19, "display", "ink", -0.4)])
    D.text(s, x + IN(0.3), IN(5.4), int(sw - IN(0.6)), IN(0.4),
           [(b, 19, "display", "accent2" if "?" in b else "mute", -0.4)])
foot(s, 14)

# ═════════════════════════════════════════════════════════════════════════════
# 15 · BÀI 5 — BÀI GIẢI CÁCH 1
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Bài giải cách 1: tính số hộp trước rồi tính số thùng.")
D.orb(s, IN(-2.4), IN(-1.6), IN(6.2), "accent2", 12)
lesson(s, "Cách 1", "accent2")
D.heading(s, "Bài giải")
D.card(s, M, IN(2.75), int(CW * 0.62), IN(3.4), "card", "accent2")
gy = IN(3.1)
for txt, size, fnt, col in [
        ("Số hộp xếp được là:", 22, "body", "ink"),
        ("800 : 4 = 200 (hộp)", 34, "display", "accent"),
        ("Số thùng bánh xếp được là:", 22, "body", "ink"),
        ("200 : 5 = 40 (thùng)", 34, "display", "accent"),
        ("Đáp số:  40 thùng", 28, "display", "accent2")]:
    D.text(s, M + IN(0.6), gy, int(CW * 0.62 - IN(1.2)), IN(0.55),
           [(txt, size, fnt, col, -0.4)])
    gy += IN(0.62) if fnt == "display" else IN(0.48)
bx = int(M + CW * 0.62 + IN(0.35))
D.card(s, bx, IN(2.75), int(CW * 0.38 - IN(0.35)), IN(3.4), "cardsoft", "stroke")
D.text(s, bx + IN(0.5), IN(3.05), IN(3.5), IN(0.4),
       [("TÓM TẮT", 12, "bold", "dim", 2.0)])
for j, (a, b) in enumerate([("Có", "800 cái bánh"), ("1 hộp", "4 cái bánh"),
                            ("1 thùng", "5 hộp"), ("Có", "…… thùng ?")]):
    y = IN(3.55) + int(IN(0.62) * j)
    D.text(s, bx + IN(0.5), y, IN(1.5), IN(0.4), [(a + ":", 17, "bold", "mute")])
    D.text(s, bx + IN(1.85), y, IN(2.4), IN(0.4),
           [(b, 17, "display", "accent2" if "?" in b else "ink", -0.3)])
foot(s, 15)

# ═════════════════════════════════════════════════════════════════════════════
# 16 · BÀI 5 — CÁCH 2 VÀ VIẾT GỌN
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Cách 2 và hai cách viết gọn của bài 5.")
D.orb(s, IN(10.6), IN(-2.0), IN(6.0), "coral", 12)
lesson(s, "Cách 2", "coral")
D.heading(s, "Còn cách nào khác không?")
D.card(s, M, IN(2.8), int(CW * 0.5 - IN(0.2)), IN(3.3), "card", "coral")
gy = IN(3.15)
for txt, size, fnt, col in [
        ("Số bánh mỗi thùng có là:", 21, "body", "ink"),
        ("4 x 5 = 20 (bánh)", 32, "display", "coral"),
        ("Số thùng bánh xếp được là:", 21, "body", "ink"),
        ("800 : 20 = 40 (thùng)", 32, "display", "coral"),
        ("Đáp số:  40 thùng", 26, "display", "accent2")]:
    D.text(s, M + IN(0.55), gy, int(CW * 0.5 - IN(1.3)), IN(0.55),
           [(txt, size, fnt, col, -0.4)])
    gy += IN(0.6) if fnt == "display" else IN(0.46)
bx = int(M + CW * 0.5 + IN(0.2))
bw = int(CW * 0.5 - IN(0.2))
D.text(s, bx, IN(2.85), bw, IN(0.35), [("VIẾT GỌN", 12, "bold", "dim", 2.0)])
for j, e in enumerate(["800 : 4 : 5 = 40 (thùng)", "800 : (4 x 5) = 40 (thùng)"]):
    y = IN(3.3) + int(IN(1.05) * j)
    D.card(s, bx, y, bw, IN(0.88), "cardmint", "accent2")
    D.text(s, bx + IN(0.45), y + IN(0.2), bw - IN(0.9), IN(0.5),
           [(e, 26, "display", "ink", -0.5)])
D.card(s, bx, IN(5.4), bw, IN(0.7), "cardwarm", "gold")
D.text(s, bx + IN(0.45), IN(5.56), bw - IN(0.9), IN(0.4),
       [("Cả hai cách đều cho đáp số 40 thùng.", 17, "bold", "ink")])
foot(s, 16)

# ═════════════════════════════════════════════════════════════════════════════
# 17 · TRÒ CHƠI
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Trò chơi Hái hoa dân chủ: ba bông hoa là ba câu hỏi ôn quy tắc.")
D.orb(s, IN(3.0), IN(0.4), IN(7.8), "petal", 18)
D.kicker(s, "Trò chơi", color="petal")
chip(s, "Trò chơi", "petal")
D.text(s, M, IN(1.9), CW, IN(1.2),
       [("Hái hoa dân chủ", 56, "display", "ink", -1.5)], name="HEADING")
qs = [("Câu 1", "Biểu thức chỉ có cộng, trừ\nhoặc chỉ có nhân, chia\nthì làm thế nào?",
       "petal"),
      ("Câu 2", "Biểu thức có cả cộng, trừ,\nnhân, chia thì làm thế nào?", "gold"),
      ("Câu 3", "Biểu thức có dấu ngoặc ( )\nthì làm thế nào?", "coral")]
qw = (CW - IN(0.5)) / 3
for i, (lab, q, col) in enumerate(qs):
    x = int(M + i * (qw + IN(0.25)))
    D.card(s, x, IN(3.35), int(qw), IN(2.75), "card", col)
    flower(s, x + int(qw / 2), IN(3.35), IN(0.42), col)
    D.text(s, x, IN(3.95), int(qw), IN(0.4), [(lab, 20, "display", col, 0.5)],
           align=PP_ALIGN.CENTER)
    D.text(s, x + IN(0.35), IN(4.5), int(qw - IN(0.7)), IN(1.4),
           [(q.replace("\n", " "), 17, "body", "ink")], line=1.35, align=PP_ALIGN.CENTER)
foot(s, 17)

# ═════════════════════════════════════════════════════════════════════════════
# 18 · BA QUY TẮC
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Đáp án trò chơi: ba quy tắc tính giá trị biểu thức.")
D.orb(s, IN(10.2), IN(4.2), IN(5.8), "accent", 14)
D.kicker(s, "Củng cố", color="accent")
chip(s, "Ghi nhớ", "accent2")
D.heading(s, "Ba quy tắc phải nhớ")
rules = [("Chỉ cộng trừ\nhoặc chỉ nhân chia", "Làm từ TRÁI sang PHẢI",
          "67 – 27 + 10 = 40 + 10 = 50", "petal"),
         ("Có cả cộng trừ\nvà nhân chia", "NHÂN CHIA trước, CỘNG TRỪ sau",
          "15 + 7 x 8 = 15 + 56 = 71", "gold"),
         ("Có dấu ngoặc ( )", "Làm TRONG NGOẶC trước",
          "123 x (42 – 40) = 123 x 2 = 246", "coral")]
rw = (CW - IN(0.5)) / 3
for i, (t, r, ex, col) in enumerate(rules):
    x = int(M + i * (rw + IN(0.25)))
    D.card(s, x, IN(2.8), int(rw), IN(3.4), "card", col)
    n = s.shapes.add_shape(MSO_SHAPE.OVAL, x + IN(0.4), IN(3.1), IN(0.56), IN(0.56))
    D.solid(n, col)
    D._no_line(n)
    D.text(s, x + IN(0.4), IN(3.22), IN(0.56), IN(0.4),
           [(str(i + 1), 22, "display", "bg2")], align=PP_ALIGN.CENTER)
    D.text(s, x + IN(0.4), IN(3.9), int(rw - IN(0.8)), IN(0.7),
           [(t.replace("\n", " "), 18, "display", "ink", -0.3)], line=1.15)
    D.text(s, x + IN(0.4), IN(4.62), int(rw - IN(0.8)), IN(0.5),
           [(r, 15, "bold", col)], line=1.25)
    D.rule(s, x + IN(0.4), IN(5.42), IN(0.5), "stroke", Pt(2))
    D.text(s, x + IN(0.4), IN(5.62), int(rw - IN(0.8)), IN(0.5),
           [(ex, 14, "num", "mute")], line=1.25)
foot(s, 18)

# ═════════════════════════════════════════════════════════════════════════════
# 19 · CHÀO TẠM BIỆT
# ═════════════════════════════════════════════════════════════════════════════
s = D.slide("Kết thúc tiết học.")
D.orb(s, IN(2.6), IN(0.2), IN(8.4), "accent", 20)
D.orb(s, IN(10.4), IN(4.4), IN(5.2), "gold", 14, "ORB2")
D.text(s, M, IN(1.5), CW, IN(0.4),
       [("TRƯỜNG TIỂU HỌC TÂN BÌNH  ·  LỚP 3.10", 13, "bold", "accent", 2.6)],
       name="KICKER")
D.rule(s, M, IN(2.06), IN(1.6), "coral", Pt(5), "RULE")
D.text(s, M, IN(2.5), IN(11.6), IN(2.0),
       [("Chân thành cảm ơn", 54, "display", "ink", -1.5),
        ("quý thầy cô và các em!", 54, "display", "accent", -1.5)],
       line=1.12, name="HEADING")
D.rule(s, M, IN(5.3), IN(11.6), "stroke", Pt(1))
D.text(s, M, IN(5.6), IN(11.6), IN(0.35),
       [("GIÁO VIÊN THỰC HIỆN", 10.5, "bold", "dim", 1.8)])
D.text(s, M, IN(5.9), IN(11.6), IN(0.45),
       [("Nguyễn Thị Thanh Tâm", 22, "display", "ink", -0.5)])
for i in range(4):
    flower(s, IN(9.5 + i * 0.9), IN(5.9), IN(0.28),
           ["petal", "gold", "coral", "grape"][i])

# ─────────────────────────────────────────────────────────────────────────────
outdir = r"F:\PPT"
out = os.path.join(outdir, "Toan3-Luyen-tap-chung-NguyenThiThanhTam.pptx")
D.save(out)
print("Da luu:", out)
