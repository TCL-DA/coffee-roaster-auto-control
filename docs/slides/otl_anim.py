# -*- coding: utf-8 -*-
"""
otl_anim.py — Bơm hiệu ứng động thật (p:timing) vào file .pptx do python-pptx sinh ra.

python-pptx KHÔNG hỗ trợ animation, nên phải ghi thẳng XML theo lược đồ
PresentationML. Module này dựng đúng cây thời gian mà PowerPoint chờ đợi:

    p:timing > p:tnLst > par(tmRoot) > seq(mainSeq) > [khối bấm] > [hiệu ứng]

Cách dùng:
    from otl_anim import animate
    animate("deck.pptx", skip_prefix="CHROME")

Nguyên tắc dàn dựng — giữ sân khấu sạch, không rối mắt:
  - Khối tiêu đề tự chạy ngay khi sang slide, người trình bày không phải bấm.
  - Phần thân chia tối đa 3 khối, mỗi khối một lần bấm để giữ nhịp nói.
  - Trong một khối: xếp theo hàng từ trên xuống, trong hàng chạy từ trái sang phải.
  - Hình trang trí (tên bắt đầu bằng CHROME) không bao giờ được gán hiệu ứng.
"""

from lxml import etree

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
EMU_IN = 914400

# Hai hình lệch nhau dưới mức này coi như cùng một hàng
ROW_TOL = int(0.30 * EMU_IN)
# Số lần bấm tối đa cho phần thân một slide
MAX_BLOCKS = 3

# presetID/presetSubtype để PowerPoint hiện đúng tên hiệu ứng trong ngăn Animation
PRESETS = {
    "rise": (10, 0),      # Fade (kèm trượt lên)
    "fade": (10, 0),      # Fade
    "wipeL": (22, 8),     # Wipe from left
    "wipeU": (22, 4),     # Wipe from bottom
    "zoom": (23, 16),     # Zoom
}


# ─────────────────────────────────────────────────────────────────────────────
# Mảnh XML của từng loại hiệu ứng
# ─────────────────────────────────────────────────────────────────────────────
def _set_visible(spid, cid):
    """Bật hiển thị tại mốc 0. Thiếu mảnh này thì hình hiện sẵn từ đầu slide."""
    return (
        '<p:set xmlns:p="%s">'
        '<p:cBhvr><p:cTn id="%d" dur="1" fill="hold">'
        '<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
        '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl>'
        '<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
        '</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>' % (NS_P, cid, spid)
    )


def _fade(spid, cid, dur):
    return (
        '<p:animEffect xmlns:p="%s" transition="in" filter="fade">'
        '<p:cBhvr><p:cTn id="%d" dur="%d"/>'
        '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl></p:cBhvr></p:animEffect>'
        % (NS_P, cid, dur, spid)
    )


def _wipe(spid, cid, dur, direction):
    return (
        '<p:animEffect xmlns:p="%s" transition="in" filter="wipe(%s)">'
        '<p:cBhvr><p:cTn id="%d" dur="%d"/>'
        '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl></p:cBhvr></p:animEffect>'
        % (NS_P, direction, cid, dur, spid)
    )


def _drift(spid, cid, dur, axis="ppt_y", amount=0.035):
    """Trượt nhẹ về đúng chỗ — cộng với fade cho cảm giác 'nâng lên' kiểu Keynote."""
    return (
        '<p:anim xmlns:p="%s" calcmode="lin" valueType="num">'
        '<p:cBhvr additive="base">'
        '<p:cTn id="%d" dur="%d" fill="hold"/>'
        '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl>'
        '<p:attrNameLst><p:attrName>%s</p:attrName></p:attrNameLst>'
        '</p:cBhvr><p:tavLst>'
        '<p:tav tm="0"><p:val><p:strVal val="#%s%+.3f"/></p:val></p:tav>'
        '<p:tav tm="100000"><p:val><p:strVal val="#%s"/></p:val></p:tav>'
        '</p:tavLst></p:anim>' % (NS_P, cid, dur, spid, axis, axis, amount, axis)
    )


def _scale(spid, cid, dur, start=0.94):
    return (
        '<p:animScale xmlns:p="%s"><p:cBhvr>'
        '<p:cTn id="%d" dur="%d" fill="hold"/>'
        '<p:tgtEl><p:spTgt spid="%d"/></p:tgtEl></p:cBhvr>'
        '<p:from x="%d" y="%d"/><p:to x="100000" y="100000"/></p:animScale>'
        % (NS_P, cid, dur, spid, int(start * 100000), int(start * 100000))
    )


def _effect(spid, kind, dur, delay, ids, node_type):
    """Một hiệu ứng vào cho một hình. `ids` là bộ đếm id dùng chung trong slide."""
    preset, subtype = PRESETS[kind]
    root_id = next(ids)
    body = [_set_visible(spid, next(ids))]

    if kind == "rise":
        body.append(_fade(spid, next(ids), dur))
        body.append(_drift(spid, next(ids), dur, "ppt_y", 0.035))
    elif kind == "fade":
        body.append(_fade(spid, next(ids), dur))
    elif kind == "wipeL":
        body.append(_wipe(spid, next(ids), dur, "left"))
    elif kind == "wipeU":
        body.append(_wipe(spid, next(ids), dur, "up"))
    elif kind == "zoom":
        body.append(_fade(spid, next(ids), dur))
        body.append(_scale(spid, next(ids), dur, 0.94))

    return (
        '<p:par xmlns:p="%s"><p:cTn id="%d" presetID="%d" presetClass="entr" '
        'presetSubtype="%d" fill="hold" grpId="0" nodeType="%s">'
        '<p:stCondLst><p:cond delay="%d"/></p:stCondLst>'
        '<p:childTnLst>%s</p:childTnLst></p:cTn></p:par>'
        % (NS_P, root_id, preset, subtype, node_type, delay, "".join(body))
    )


def _block(effects_xml, ids, on_click):
    """Một lần bấm chuột = một khối; bên trong, các hiệu ứng tự chạy theo delay."""
    outer, inner = next(ids), next(ids)
    start = '<p:cond delay="indefinite"/>' if on_click else '<p:cond delay="0"/>'
    return (
        '<p:par xmlns:p="%s"><p:cTn id="%d" fill="hold">'
        '<p:stCondLst>%s</p:stCondLst><p:childTnLst>'
        '<p:par><p:cTn id="%d" fill="hold">'
        '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
        '<p:childTnLst>%s</p:childTnLst></p:cTn></p:par>'
        '</p:childTnLst></p:cTn></p:par>'
        % (NS_P, outer, start, inner, "".join(effects_xml))
    )


def _timing(blocks_xml):
    return (
        '<p:timing xmlns:p="%s"><p:tnLst>'
        '<p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
        '<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        '<p:childTnLst>%s</p:childTnLst></p:cTn>'
        '<p:prevCondLst><p:cond evt="onPrev" delay="0">'
        '<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        '<p:nextCondLst><p:cond evt="onNext" delay="0">'
        '<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        '</p:seq></p:childTnLst></p:cTn></p:par>'
        '</p:tnLst></p:timing>' % (NS_P, "".join(blocks_xml))
    )


# ─────────────────────────────────────────────────────────────────────────────
# Dàn dựng
# ─────────────────────────────────────────────────────────────────────────────
def _counter(start=3):
    n = start
    while True:
        yield n
        n += 1


def _rows(shapes):
    """Gom hình thành hàng theo toạ độ dọc; trong hàng xếp trái → phải."""
    shapes = sorted(shapes, key=lambda s: (s["top"], s["left"]))
    rows, cur, base = [], [], None
    for sh in shapes:
        if base is None:
            base, cur = sh["top"], [sh]
        elif sh["top"] - base <= ROW_TOL:
            cur.append(sh)
        else:
            rows.append(sorted(cur, key=lambda s: s["left"]))
            cur, base = [sh], sh["top"]
    if cur:
        rows.append(sorted(cur, key=lambda s: s["left"]))
    return rows


def _split_blocks(rows, max_blocks=MAX_BLOCKS):
    """Cắt danh sách hàng thành các khối bấm, ưu tiên cắt ở khe dọc lớn nhất."""
    if len(rows) <= 1:
        return [rows] if rows else []
    gaps = []
    for i in range(1, len(rows)):
        prev_bottom = max(r["top"] + r["height"] for r in rows[i - 1])
        gaps.append((rows[i][0]["top"] - prev_bottom, i))
    gaps.sort(key=lambda g: g[0], reverse=True)
    cuts = sorted(i for _, i in gaps[: max_blocks - 1])
    blocks, prev = [], 0
    for c in cuts + [len(rows)]:
        if c > prev:
            blocks.append(rows[prev:c])
            prev = c
    return blocks


# Tiêu đề trang bìa / trang kết: chạy ngay khi sang slide, không bắt bấm
HEADER_HINTS = ("HEADING", "TITLE_RULE")


def _is_header(name):
    return any(h in name for h in HEADER_HINTS)


def _plan_slide(shapes):
    """Trả về danh sách khối [(on_click, [(spid, kind, dur, delay), ...]), ...]."""
    header = sorted((s for s in shapes if _is_header(s["name"])),
                    key=lambda s: s["top"])
    body = [s for s in shapes if not _is_header(s["name"])]

    plan = []
    if header:
        eff, t = [], 0
        for s in header:
            eff.append((s["id"], "rise", 520, t))
            t += 110
        plan.append((False, eff))   # tự chạy khi vừa sang slide

    for block in _split_blocks(_rows(body)):
        eff, t = [], 0
        for row in block:
            wide_row = len(row) >= 3
            for s in row:
                # Hàng nhiều phần tử nằm ngang → quét ngang cho có hướng đọc
                kind = "wipeL" if wide_row and s["width"] > s["height"] * 1.6 else "rise"
                eff.append((s["id"], kind, 520 if kind == "wipeL" else 460, t))
                t += 70 if wide_row else 90
            t += 90
        plan.append((True, eff))    # mỗi khối thân chờ một lần bấm
    return plan


def animate(pptx_path, skip_prefix="CHROME", verbose=True):
    """Ghi p:timing vào mọi slide của file .pptx (sửa tại chỗ). Trả về số hiệu ứng."""
    from pptx import Presentation

    pptx_path = str(pptx_path)
    prs = Presentation(pptx_path)
    total, clicks = 0, 0

    for n, slide in enumerate(prs.slides, 1):
        shapes = []
        for sh in slide.shapes:
            if sh.name.startswith(skip_prefix):
                continue
            # python-pptx đặt tên trùng nhau giữa các slide ("TextBox 20"...).
            # Morph byObject sẽ ghép nhầm rồi bay hình lung tung, chồng lên
            # hiệu ứng vào. Gắn số slide vào tên để mỗi hình là duy nhất.
            sh.name = "S%02d_%s" % (n, sh.name.replace(" ", "_"))
            if sh.left is None or sh.top is None:
                continue
            shapes.append({
                "id": sh.shape_id, "name": sh.name,
                "left": sh.left, "top": sh.top,
                "width": sh.width or 0, "height": sh.height or 0,
            })
        plan = _plan_slide(shapes)
        if not plan:
            continue

        ids = _counter(3)
        blocks = []
        for on_click, effects in plan:
            xml = []
            for i, (spid, kind, dur, delay) in enumerate(effects):
                node = ("clickEffect" if on_click else "afterEffect") if i == 0 else "withEffect"
                xml.append(_effect(spid, kind, dur, delay, ids, node))
                total += 1
            blocks.append(_block(xml, ids, on_click))
            clicks += 1 if on_click else 0

        sld = slide._element
        for old in sld.findall("{%s}timing" % NS_P):
            sld.remove(old)
        sld.append(etree.fromstring(_timing(blocks)))

    prs.save(pptx_path)
    if verbose:
        print(f"Animation: {total} hieu ung / {clicks} lan bam -> {pptx_path}")
    return total
